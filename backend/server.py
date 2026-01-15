from fastapi import FastAPI, APIRouter, File, UploadFile, Form, HTTPException
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from contextlib import asynccontextmanager
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware as StarletteCORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional
import uuid
from datetime import datetime, timezone
import io
import zipfile
import shutil
from PIL import Image
from pypdf import PdfReader, PdfWriter
from pptx import Presentation
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, Border, Side, Alignment
import pdfplumber
from xlsx2pdf.transformator import Transformer

# ============== FILE SIZE CONSTANTS ==============
# Maximum file size: 30MB (30 * 1024 * 1024 bytes)
MAX_FILE_SIZE = 30 * 1024 * 1024  # 30 MB in bytes

# Helper function to format file size for display
def format_file_size(bytes_size):
    """Format file size in human-readable format"""
    if bytes_size == 0:
        return '0 Bytes'
    units = ['Bytes', 'KB', 'MB', 'GB']
    i = 0
    while bytes_size >= 1024 and i < len(units) - 1:
        bytes_size /= 1024
        i += 1
    return f"{bytes_size:.2f} {units[i]}"

# Import watermark service
from services.watermark_service import (
    add_text_watermark,
    add_image_watermark,
    add_multiple_watermarks,
    POSITION_CENTER,
    POSITION_TOP_LEFT,
    POSITION_TOP_RIGHT,
    POSITION_BOTTOM_LEFT,
    POSITION_BOTTOM_RIGHT,
    POSITION_TILED
)

# Import table extraction service for precise PDF to Excel conversion
from services.table_extraction_service import (
    TableExtractionService,
    TableData,
    TableInfo,
    TableExtractionMethod
)


def find_font_path() -> str:
    """Find a suitable TrueType font for xlsx2pdf conversion.
    
    The Transformer class requires a font file path. This function searches
    common font locations on Linux systems.
    
    Returns:
        str: Path to a valid font file
        
    Raises:
        FileNotFoundError: If no suitable font is found
    """
    # Common Linux font paths
    font_candidates = [
        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
        '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
        '/usr/share/fonts/truetype/msttcorefonts/Arial.ttf',
        '/usr/share/fonts/truetype/msttcorefonts/Arial_Bold.ttf',
        '/usr/share/fonts/truetype/freefont/FreeSans.ttf',
        '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
        '/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf',
        '/usr/share/fonts/TTF/DejaVuSans.ttf',
        '/usr/share/fonts/dejavu/DejaVuSans.ttf',
        '/System/Library/Fonts/Arial.ttf',  # macOS fallback
    ]
    
    for font_path in font_candidates:
        if Path(font_path).exists():
            print(f"Using font: {font_path}")
            return font_path
    
    # Try to find any .ttf file in common font directories
    import glob
    font_dirs = [
        '/usr/share/fonts/**/*.ttf',
        '/usr/local/share/fonts/**/*.ttf',
    ]
    
    for pattern in font_dirs:
        fonts = glob.glob(pattern, recursive=True)
        if fonts:
            # Return the first found font
            font_path = fonts[0]
            print(f"Using discovered font: {font_path}")
            return font_path
    
    raise FileNotFoundError(
        "No suitable TrueType font found for xlsx2pdf conversion. "
        "Please install a font package (e.g., fonts-dejavu-core) or "
        "specify a font path manually."
    )
from pdf2docx import Converter
from docx2pdf import convert as docx2pdf_convert
import pytesseract
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
import tempfile
from pathlib import Path
import uuid
from pypdf import PdfReader
import subprocess
import shutil as tesseract_shutil

# Import document conversion routes
# All document conversion routes are defined directly in this file

# Configure Tesseract executable path
# Uncomment and set the path if Tesseract is not in PATH
# pytesseract.pytesseract.tesseract_cmd = r'/usr/bin/tesseract'  # Linux
# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  # Windows

# Try to auto-detect Tesseract path
def find_tesseract_executable():
    """Find Tesseract executable path"""
    # Common Linux paths
    linux_paths = ['/usr/bin/tesseract', '/usr/local/bin/tesseract']
    # Check if tesseract is in PATH
    tesseract_path = tesseract_shutil.which('tesseract')
    if tesseract_path:
        return tesseract_path
    # Check common Linux paths
    for path in linux_paths:
        if Path(path).exists():
            return path
    return None

TESSERACT_PATH = find_tesseract_executable()
if TESSERACT_PATH:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH
    print(f"Tesseract found at: {TESSERACT_PATH}")
else:
    print("Tesseract not found in PATH or common locations. OCR may not work.")

def get_tesseract_languages():
    """Get list of installed Tesseract language packs"""
    try:
        if TESSERACT_PATH:
            result = subprocess.run(
                [TESSERACT_PATH, '--list-langs'],
                capture_output=True,
                text=True,
                timeout=10
            )
            if result.returncode == 0:
                # Parse languages from output (skip first line which is "List of available languages:")
                lines = result.stdout.strip().split('\n')
                if len(lines) > 1:
                    return [lang.strip() for lang in lines[1:] if lang.strip()]
        # Fallback: try with pytesseract
        return pytesseract.get_languages(config='')
    except Exception as e:
        print(f"Failed to get Tesseract languages: {e}")
        return ['eng']  # Return at least English as fallback

# Pre-load available languages at startup
AVAILABLE_OCR_LANGUAGES = []
try:
    AVAILABLE_OCR_LANGUAGES = get_tesseract_languages()
    print(f"Available OCR languages: {AVAILABLE_OCR_LANGUAGES}")
except Exception as e:
    print(f"Could not load OCR languages: {e}")
    AVAILABLE_OCR_LANGUAGES = ['eng']

# Define common language names for display
LANGUAGE_NAMES = {
    # Script languages
    'Arabic': 'Arabic',
    'Armenian': 'Armenian',
    'Bengali': 'Bengali',
    'Canadian_Aboriginal': 'Canadian Aboriginal',
    'Cherokee': 'Cherokee',
    'Cyrillic': 'Cyrillic',
    'Devanagari': 'Devanagari',
    'Ethiopic': 'Ethiopic',
    'Fraktur': 'Fraktur',
    'Georgian': 'Georgian',
    'Greek': 'Greek',
    'Gujarati': 'Gujarati',
    'Gurmukhi': 'Gurmukhi',
    'HanS': 'Chinese Simplified',
    'HanS_vert': 'Chinese Simplified (Vertical)',
    'HanT': 'Chinese Traditional',
    'HanT_vert': 'Chinese Traditional (Vertical)',
    'Hangul': 'Hangul (Korean)',
    'Hangul_vert': 'Hangul (Korean, Vertical)',
    'Hebrew': 'Hebrew',
    'Japanese': 'Japanese',
    'Japanese_vert': 'Japanese (Vertical)',
    'Kannada': 'Kannada',
    'Khmer': 'Khmer',
    'Lao': 'Lao',
    'Latin': 'Latin',
    'Malayalam': 'Malayalam',
    'Myanmar': 'Myanmar',
    'Oriya': 'Oriya',
    'Sinhala': 'Sinhala',
    'Syriac': 'Syriac',
    'Tamil': 'Tamil',
    'Telugu': 'Telugu',
    'Thaana': 'Thaana',
    'Thai': 'Thai',
    'Tibetan': 'Tibetan',
    'Vietnamese': 'Vietnamese',
    
    # ISO 639 language codes
    'afr': 'Afrikaans',
    'amh': 'Amharic',
    'ara': 'Arabic',
    'asm': 'Assamese',
    'aze': 'Azerbaijani',
    'aze_cyrl': 'Azerbaijani (Cyrillic)',
    'bel': 'Belarusian',
    'ben': 'Bengali',
    'bod': 'Tibetan',
    'bos': 'Bosnian',
    'bre': 'Breton',
    'bul': 'Bulgarian',
    'cat': 'Catalan',
    'ceb': 'Cebuano',
    'ces': 'Czech',
    'chi_sim': 'Chinese (Simplified)',
    'chi_sim_vert': 'Chinese Simplified (Vertical)',
    'chi_tra': 'Chinese (Traditional)',
    'chi_tra_vert': 'Chinese Traditional (Vertical)',
    'chr': 'Cherokee',
    'cos': 'Corsican',
    'cym': 'Welsh',
    'dan': 'Danish',
    'deu': 'German',
    'div': 'Dhivehi',
    'dzo': 'Dzongkha',
    'ell': 'Greek',
    'eng': 'English',
    'enm': 'English (Middle)',
    'epo': 'Esperanto',
    'est': 'Estonian',
    'eus': 'Basque',
    'fao': 'Faroese',
    'fas': 'Persian',
    'fil': 'Filipino',
    'fin': 'Finnish',
    'fra': 'French',
    'frk': 'German (Fraktur)',
    'frm': 'French (Middle)',
    'fry': 'Frisian',
    'gla': 'Scottish Gaelic',
    'gle': 'Irish',
    'glg': 'Galician',
    'grc': 'Greek (Ancient)',
    'guj': 'Gujarati',
    'hat': 'Haitian Creole',
    'heb': 'Hebrew',
    'hin': 'Hindi',
    'hrv': 'Croatian',
    'hun': 'Hungarian',
    'hye': 'Armenian',
    'iku': 'Inuktitut',
    'ind': 'Indonesian',
    'isl': 'Icelandic',
    'ita': 'Italian',
    'ita_old': 'Italian (Old)',
    'jav': 'Javanese',
    'jpn': 'Japanese',
    'jpn_vert': 'Japanese (Vertical)',
    'kan': 'Kannada',
    'kat': 'Georgian',
    'kat_old': 'Georgian (Old)',
    'kaz': 'Kazakh',
    'khm': 'Khmer',
    'kir': 'Kyrgyz',
    'kmr': 'Kurmanji Kurdish',
    'kor': 'Korean',
    'kor_vert': 'Korean (Vertical)',
    'lao': 'Lao',
    'lat': 'Latin',
    'lav': 'Latvian',
    'lit': 'Lithuanian',
    'ltz': 'Luxembourgish',
    'mal': 'Malayalam',
    'mar': 'Marathi',
    'mkd': 'Macedonian',
    'mlt': 'Maltese',
    'mon': 'Mongolian',
    'mri': 'Maori',
    'msa': 'Malay',
    'mya': 'Burmese',
    'nep': 'Nepali',
    'nld': 'Dutch',
    'nor': 'Norwegian',
    'oci': 'Occitan',
    'ori': 'Oriya',
    'osd': 'Orientation and Script Detection',
    'pan': 'Punjabi',
    'pol': 'Polish',
    'por': 'Portuguese',
    'pus': 'Pashto',
    'que': 'Quechua',
    'ron': 'Romanian',
    'rus': 'Russian',
    'san': 'Sanskrit',
    'sin': 'Sinhala',
    'slk': 'Slovak',
    'slv': 'Slovenian',
    'snd': 'Sindhi',
    'spa': 'Spanish',
    'spa_old': 'Spanish (Old)',
    'sqi': 'Albanian',
    'srp': 'Serbian',
    'srp_latn': 'Serbian (Latin)',
    'sun': 'Sundanese',
    'swa': 'Swahili',
    'swe': 'Swedish',
    'syr': 'Syriac',
    'tam': 'Tamil',
    'tat': 'Tatar',
    'tel': 'Telugu',
    'tgk': 'Tajik',
    'tha': 'Thai',
    'tir': 'Tigrinya',
    'ton': 'Tongan',
    'tur': 'Turkish',
    'uig': 'Uyghur',
    'ukr': 'Ukrainian',
    'urd': 'Urdu',
    'uzb': 'Uzbek',
    'uzb_cyrl': 'Uzbek (Cyrillic)',
    'vie': 'Vietnamese',
    'yid': 'Yiddish',
    'yor': 'Yoruba',
}

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Define lifespan context manager for startup and shutdown events
@asynccontextmanager
async def lifespan(app: FastAPI):
    """Lifespan context manager for startup and shutdown events."""
    # Startup: app is starting
    print("Starting up File Conversion API...")
    yield
    # Shutdown: app is closing
    print("Shutting down File Conversion API...")
    client.close()

# Create the main app with lifespan
app = FastAPI(lifespan=lifespan)

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")

# Create temp directory for file operations
# TEMP_DIR = Path("/tmp/file_conversions")
# TEMP_DIR.mkdir(exist_ok=True)
# make TEMP_DIR Windows-friendly and auto-create
TEMP_DIR = Path(os.getenv("TEMP_DIR", Path.cwd() / "tmp" / "file_conversions"))
TEMP_DIR.mkdir(parents=True, exist_ok=True)


# Define Models
class ConversionHistory(BaseModel):
    model_config = ConfigDict(extra="ignore")
    
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    conversion_type: str
    source_format: str
    target_format: str
    filename: str
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    status: str

class ConversionHistoryCreate(BaseModel):
    conversion_type: str
    source_format: str
    target_format: str
    filename: str
    status: str

# Helper Functions
def save_upload_file_tmp(upload_file: UploadFile) -> Path:
    """Save uploaded file to temp directory"""
    try:
        suffix = Path(upload_file.filename).suffix
        tmp_path = TEMP_DIR / f"{uuid.uuid4()}{suffix}"
        with tmp_path.open("wb") as buffer:
            shutil.copyfileobj(upload_file.file, buffer)
        return tmp_path
    finally:
        upload_file.file.close()

def convert_image_format(input_path: Path, output_format: str) -> Path:
    """Convert image to different format"""
    img = Image.open(input_path)
    if img.mode == 'RGBA' and output_format.lower() in ['jpg', 'jpeg']:
        img = img.convert('RGB')
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}.{output_format.lower()}"
    img.save(output_path, format=output_format.upper())
    return output_path


# ============== Image Resize Functions ==============

# Resize quality presets (correspond to PIL quality settings)
RESIZE_QUALITY_PRESETS = {
    'low': 50,       # 50% quality - smaller file size
    'medium': 75,    # 75% quality - balanced
    'high': 90,      # 90% quality - good quality
    'maximum': 100,  # 100% quality - best quality
}


def resize_image(
    image_path: Path,
    target_width: int,
    target_height: int,
    maintain_aspect_ratio: bool = True,
    output_format: str = 'jpeg',
    quality: str = 'high'
) -> Path:
    """Resize an image to specified dimensions.
    
    Args:
        image_path: Path to the input image
        target_width: Target width in pixels
        target_height: Target height in pixels
        maintain_aspect_ratio: If True, maintains aspect ratio (default: True)
        output_format: Output format (jpeg, png, webp, bmp)
        quality: Quality preset (low, medium, high, maximum)
    
    Returns:
        Path to the resized image
    """
    img = Image.open(image_path)
    
    # Get original dimensions
    orig_width, orig_height = img.size
    
    # Calculate new dimensions
    if maintain_aspect_ratio:
        # Calculate aspect ratio
        orig_aspect = orig_width / orig_height if orig_height > 0 else 1
        target_aspect = target_width / target_height if target_height > 0 else 1
        
        if orig_aspect > target_aspect:
            # Width is the limiting factor
            new_width = target_width
            new_height = int(target_width / orig_aspect)
        else:
            # Height is the limiting factor
            new_height = target_height
            new_width = int(target_height * orig_aspect)
    else:
        # Stretch to exact dimensions
        new_width = target_width
        new_height = target_height
    
    # Ensure minimum dimensions
    new_width = max(1, new_width)
    new_height = max(1, new_height)
    
    # Resize image using high-quality resampling
    resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
    
    # Handle transparency for JPEG format
    if output_format.lower() in ['jpg', 'jpeg'] and resized_img.mode in ('RGBA', 'LA', 'P'):
        resized_img = resized_img.convert('RGB')
    
    # Determine output format name for PIL
    pil_format = output_format.upper()
    if pil_format == 'JPG':
        pil_format = 'JPEG'
    
    # Get quality value
    quality_value = RESIZE_QUALITY_PRESETS.get(quality, RESIZE_QUALITY_PRESETS['high'])
    
    # Save resized image
    output_path = TEMP_DIR / f"{uuid.uuid4()}.{output_format.lower()}"
    
    # For PNG, use compression level; for others use quality
    if output_format.lower() == 'png':
        resized_img.save(output_path, format=pil_format, optimize=True)
    else:
        resized_img.save(output_path, format=pil_format, quality=quality_value)
    
    return output_path


def resize_multiple_images(
    image_paths: List[Path],
    target_width: int,
    target_height: int,
    maintain_aspect_ratio: bool = True,
    output_format: str = 'jpeg',
    quality: str = 'high'
) -> Path:
    """Resize multiple images and return as ZIP.
    
    Args:
        image_paths: List of paths to input images
        target_width: Target width in pixels
        target_height: Target height in pixels
        maintain_aspect_ratio: If True, maintains aspect ratio
        output_format: Output format (jpeg, png, webp, bmp)
        quality: Quality preset
    
    Returns:
        Path to ZIP containing resized images
    """
    resized_paths = []
    
    for image_path in image_paths:
        try:
            resized_path = resize_image(
                image_path,
                target_width,
                target_height,
                maintain_aspect_ratio,
                output_format,
                quality
            )
            resized_paths.append(resized_path)
        except Exception as e:
            print(f"Failed to resize {image_path}: {e}")
            continue
    
    if not resized_paths:
        raise Exception("No images could be resized")
    
    # Create ZIP with resized images
    zip_name = f"resized_images_{uuid.uuid4()}"
    zip_path = create_zip(resized_paths, zip_name)
    
    # Clean up individual resized images
    for path in resized_paths:
        try:
            path.unlink()
        except:
            pass
    
    return zip_path


def lock_pdf(pdf_path: Path, password: str) -> Path:
    """Encrypt PDF with password"""
    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    
    for page in reader.pages:
        writer.add_page(page)
    
    writer.encrypt(password)
    output_path = TEMP_DIR / f"{uuid.uuid4()}_locked.pdf"
    with open(output_path, "wb") as output_file:
        writer.write(output_file)
    
    return output_path

def unlock_pdf(pdf_path: Path, password: str) -> Path:
    """Decrypt PDF with password"""
    reader = PdfReader(pdf_path)
    if reader.is_encrypted:
        reader.decrypt(password)
    
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}_unlocked.pdf"
    with open(output_path, "wb") as output_file:
        writer.write(output_file)
    
    return output_path

def merge_pdfs(pdf_paths: List[Path]) -> Path:
    """Merge multiple PDFs"""
    writer = PdfWriter()
    
    for pdf_path in pdf_paths:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            writer.add_page(page)
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}_merged.pdf"
    with open(output_path, "wb") as output_file:
        writer.write(output_file)
    
    return output_path

def split_pdf(pdf_path: Path, page_ranges: str) -> List[Path]:
    """Split PDF into multiple files"""
    reader = PdfReader(pdf_path)
    output_paths = []
    
    ranges = page_ranges.split(',')
    for idx, range_str in enumerate(ranges):
        writer = PdfWriter()
        if '-' in range_str:
            start, end = map(int, range_str.split('-'))
            for page_num in range(start - 1, min(end, len(reader.pages))):
                writer.add_page(reader.pages[page_num])
        else:
            page_num = int(range_str) - 1
            if page_num < len(reader.pages):
                writer.add_page(reader.pages[page_num])
        
        output_path = TEMP_DIR / f"{uuid.uuid4()}_part{idx + 1}.pdf"
        with open(output_path, "wb") as output_file:
            writer.write(output_file)
        output_paths.append(output_path)
    
    return output_paths

def create_zip(file_paths: List[Path], zip_name: str, base_dir: Path = None) -> Path:
    """Create ZIP archive preserving folder structure"""
    output_path = TEMP_DIR / f"{uuid.uuid4()}_{zip_name}.zip"
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        if base_dir:
            # Preserve folder structure relative to base_dir
            for file_path in file_paths:
                arcname = str(file_path.relative_to(base_dir.parent))
                zipf.write(file_path, arcname)
        else:
            for file_path in file_paths:
                zipf.write(file_path, file_path.name)
    return output_path

def extract_zip(zip_path: Path) -> List[Path]:
    """Extract ZIP archive"""
    extract_dir = TEMP_DIR / f"{uuid.uuid4()}_extracted"
    extract_dir.mkdir(exist_ok=True)
    
    with zipfile.ZipFile(zip_path, 'r') as zipf:
        zipf.extractall(extract_dir)
    
    return list(extract_dir.glob('*'))

def ocr_image(image_path: Path, language: str = "eng") -> str:
    """Extract text from image using OCR"""
    try:
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img, lang=language)
        return text
    except Exception as e:
        error_msg = str(e).lower()
        if "language" in error_msg or "lang" in error_msg:
            return f"OCR Error: Language '{language}' not found. Please install the language pack or choose a different language."
        elif "tesseract" in error_msg or "not found" in error_msg:
            return "OCR Error: Tesseract is not installed or not found. Please install Tesseract OCR."
        else:
            return f"OCR Error: {str(e)}. Make sure Tesseract is installed and language pack is available."

def search_in_pdf(pdf_path: Path, search_term: str) -> dict:
    """Search for text in PDF and return results with page numbers and context"""
    reader = PdfReader(str(pdf_path))
    results = []
    search_lower = search_term.lower()

    for page_num, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and search_lower in text.lower():
            # Find all occurrences with context
            text_lower = text.lower()
            start = 0
            while True:
                pos = text_lower.find(search_lower, start)
                if pos == -1:
                    break

                # Extract context around the match (100 chars before and after)
                context_start = max(0, pos - 100)
                context_end = min(len(text), pos + len(search_term) + 100)
                context = text[context_start:context_end]

                # Highlight the search term in context
                highlighted_context = context.replace(
                    text[pos:pos + len(search_term)],
                    f"**{text[pos:pos + len(search_term)]}**"
                )

                results.append({
                    "page": page_num + 1,
                    "context": highlighted_context,
                    "position": pos
                })

                start = pos + 1

    return {
        "total_matches": len(results),
        "search_term": search_term,
        "results": results
    }

def detect_language_from_image(image_path: Path) -> dict:
    """Detect language/script from image using Tesseract OSD"""
    try:
        img = Image.open(image_path)
        
        # Use Tesseract OSD (Orientation and Script Detection)
        osd_data = pytesseract.image_to_osd(img)
        
        # Parse OSD output
        osd_lines = osd_data.split('\n')
        detected_script = None
        detected_orientation = 0
        confidence = 0
        
        for line in osd_lines:
            if 'Script:' in line:
                detected_script = line.split(':')[1].strip()
            elif 'Orientation in degrees:' in line:
                detected_orientation = int(line.split(':')[1].strip())
            elif 'Rotate:' in line:
                # Rotate value
                pass
            elif 'Confidence:' in line:
                try:
                    confidence = float(line.split(':')[1].strip())
                except:
                    pass
        
        # Map script to likely language codes
        suggested_languages = []
        script_to_languages = {
            'Latin': ['eng', 'fra', 'deu', 'spa', 'ita', 'por', 'nld', 'pol', 'ces', 'dan', 'fin', 'nor', 'swe', 'ron', 'hun', 'cat', 'glg', 'eusk', 'gle', 'bre', 'lat'],
            'Cyrillic': ['rus', 'ukr', 'bul', 'bel', 'srp', 'mkd', 'kaz', 'uzb', 'tgk', 'kir'],
            'Arabic': ['ara', 'fas', 'urd', 'pus', 'div', 'snd'],
            'Devanagari': ['hin', 'mar', 'nep', 'san', 'bod'],
            'Bengali': ['ben', 'asm'],
            'Tamil': ['tam'],
            'Telugu': ['tel'],
            'Kannada': ['kan'],
            'Malayalam': ['mal'],
            'Gujarati': ['guj'],
            'Oriya': ['ori'],
            'Punjabi': ['pan'],
            'Myanmar': ['mya'],
            'Thai': ['tha'],
            'Lao': ['lao'],
            'Khmer': ['khm'],
            'Hebrew': ['heb'],
            'Greek': ['ell'],
            'Japanese': ['jpn'],
            'Korean': ['kor'],
            'Chinese': ['chi_sim', 'chi_tra'],
            'HanS': ['chi_sim'],
            'HanT': ['chi_tra'],
            'Hangul': ['kor'],
            'Hangul_vert': ['kor'],
            'HanS_vert': ['chi_sim'],
            'HanT_vert': ['chi_tra'],
            'Japanese_vert': ['jpn'],
        }
        
        if detected_script in script_to_languages:
            suggested_languages = script_to_languages[detected_script]
        else:
            # Fallback: try to find partial match
            for script, langs in script_to_languages.items():
                if detected_script and (script.lower() in detected_script.lower() or detected_script.lower() in script.lower()):
                    suggested_languages = langs
                    break
            # If still no match, try to find matching available languages
            if not suggested_languages and detected_script:
                for lang_code, lang_name in LANGUAGE_NAMES.items():
                    if isinstance(lang_name, str) and detected_script.lower() in lang_name.lower():
                        if lang_code not in suggested_languages:
                            suggested_languages.append(lang_code)
        
        # Filter to only available languages
        available_suggestions = []
        for lang in suggested_languages:
            if lang in AVAILABLE_OCR_LANGUAGES:
                available_suggestions.append(lang)
        
        # If no suggestions, return English as fallback
        if not available_suggestions:
            if 'eng' in AVAILABLE_OCR_LANGUAGES:
                available_suggestions = ['eng']
            elif AVAILABLE_OCR_LANGUAGES:
                available_suggestions = [AVAILABLE_OCR_LANGUAGES[0]]
        
        return {
            "detected_script": detected_script,
            "orientation": detected_orientation,
            "confidence": confidence,
            "suggested_languages": available_suggestions[:5]  # Return top 5 suggestions
        }
    except Exception as e:
        print(f"Language detection failed: {e}")
        # Return fallback
        if 'eng' in AVAILABLE_OCR_LANGUAGES:
            return {
                "detected_script": "Unknown",
                "orientation": 0,
                "confidence": 0,
                "suggested_languages": ['eng']
            }
        elif AVAILABLE_OCR_LANGUAGES:
            return {
                "detected_script": "Unknown",
                "orientation": 0,
                "confidence": 0,
                "suggested_languages": [AVAILABLE_OCR_LANGUAGES[0]]
            }
        return {
            "detected_script": "Unknown",
            "orientation": 0,
            "confidence": 0,
            "suggested_languages": []
        }
    





 #============== DOCX Conversion Functions ==============

def convert_docx_to_pdf_libreoffice(docx_path: Path) -> Path:
    """Convert DOCX to PDF using LibreOffice in headless mode.
    
    This is the preferred method on Linux as it works natively without Wine.
    Command: libreoffice --headless --convert-to pdf --outdir /output /input.docx
    """
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pdf"
    
    try:
        # LibreOffice requires a HOME directory to work properly
        # Set HOME env var if not already set
        env = os.environ.copy()
        if 'HOME' not in env:
            env['HOME'] = '/tmp'
        
        # LibreOffice command: headless mode, convert to PDF, specify output directory
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(TEMP_DIR),
            str(docx_path)
        ]
        
        print(f"Running LibreOffice command: {' '.join(cmd)}")
        
        # Run the conversion
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120,  # 2 minute timeout
            env=env
        )
        
        if result.returncode != 0:
            print(f"LibreOffice stderr: {result.stderr}")
            raise Exception(f"LibreOffice conversion failed with return code {result.returncode}")
        
        # LibreOffice creates the PDF with the same name as input but .pdf extension
        # in the output directory
        expected_pdf_path = TEMP_DIR / f"{docx_path.stem}.pdf"
        
        if expected_pdf_path.exists() and expected_pdf_path.stat().st_size > 0:
            # Move to our output path with unique name
            shutil.move(str(expected_pdf_path), str(output_path))
            return output_path
        else:
            # Check if LibreOffice put it elsewhere
            # Sometimes it outputs to current directory
            fallback_pdf = Path(f"{docx_path.stem}.pdf")
            if fallback_pdf.exists() and fallback_pdf.stat().st_size > 0:
                shutil.move(str(fallback_pdf), str(output_path))
                return output_path
            
            raise Exception(f"LibreOffice conversion failed - output PDF not found at expected location")
            
    except subprocess.TimeoutExpired:
        raise Exception("LibreOffice conversion timed out after 120 seconds")
    except FileNotFoundError:
        raise Exception("LibreOffice is not installed or not found in PATH")
    except Exception as e:
        raise Exception(f"LibreOffice conversion failed: {str(e)}")


def check_libreoffice_available() -> bool:
    """Check if LibreOffice is available on the system."""
    try:
        result = subprocess.run(
            ['libreoffice', '--version'],
            capture_output=True,
            text=True,
            timeout=10
        )
        if result.returncode == 0:
            print(f"LibreOffice available: {result.stdout.strip()}")
            return True
    except (subprocess.SubprocessError, FileNotFoundError):
        pass
    return False


def convert_docx_to_doc(docx_path: Path) -> Path:
    """Convert DOCX to DOC using LibreOffice in headless mode.
    
    This uses LibreOffice which is already installed in the system.
    Command: libreoffice --headless --convert-to doc --outdir /output /input.docx
    """
    output_path = TEMP_DIR / f"{uuid.uuid4()}.doc"
    
    try:
        # LibreOffice requires a HOME directory to work properly
        env = os.environ.copy()
        if 'HOME' not in env:
            env['HOME'] = '/tmp'
        
        # LibreOffice command: headless mode, convert to DOC, specify output directory
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'doc',
            '--outdir', str(TEMP_DIR),
            str(docx_path)
        ]
        
        print(f"Running LibreOffice command: {' '.join(cmd)}")
        
        # Run the conversion
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120,
            env=env
        )
        
        if result.returncode != 0:
            print(f"LibreOffice stderr: {result.stderr}")
            raise Exception(f"LibreOffice conversion failed with return code {result.returncode}")
        
        # LibreOffice creates the DOC with the same name as input but .doc extension
        expected_doc_path = TEMP_DIR / f"{docx_path.stem}.doc"
        
        if expected_doc_path.exists() and expected_doc_path.stat().st_size > 0:
            # Move to our output path with unique name
            shutil.move(str(expected_doc_path), str(output_path))
            return output_path
        else:
            # Check if LibreOffice put it elsewhere
            fallback_doc = Path(f"{docx_path.stem}.doc")
            if fallback_doc.exists() and fallback_doc.stat().st_size > 0:
                shutil.move(str(fallback_doc), str(output_path))
                return output_path
            
            raise Exception(f"LibreOffice conversion failed - output DOC not found at expected location")
            
    except subprocess.TimeoutExpired:
        raise Exception("LibreOffice conversion timed out after 120 seconds")
    except FileNotFoundError:
        raise Exception("LibreOffice is not installed or not found in PATH")
    except Exception as e:
        raise Exception(f"LibreOffice conversion failed: {str(e)}")


def convert_doc_to_docx(doc_path: Path) -> Path:
    """Convert DOC to DOCX using LibreOffice in headless mode.
    
    This uses LibreOffice which is already installed in the system.
    Command: libreoffice --headless --convert-to docx --outdir /output /input.doc
    """
    output_path = TEMP_DIR / f"{uuid.uuid4()}.docx"
    
    try:
        # LibreOffice requires a HOME directory to work properly
        env = os.environ.copy()
        if 'HOME' not in env:
            env['HOME'] = '/tmp'
        
        # LibreOffice command: headless mode, convert to DOCX, specify output directory
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'docx',
            '--outdir', str(TEMP_DIR),
            str(doc_path)
        ]
        
        print(f"Running LibreOffice command: {' '.join(cmd)}")
        
        # Run the conversion
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120,
            env=env
        )
        
        if result.returncode != 0:
            print(f"LibreOffice stderr: {result.stderr}")
            raise Exception(f"LibreOffice conversion failed with return code {result.returncode}")
        
        # LibreOffice creates the DOCX with the same name as input but .docx extension
        expected_docx_path = TEMP_DIR / f"{doc_path.stem}.docx"
        
        if expected_docx_path.exists() and expected_docx_path.stat().st_size > 0:
            # Move to our output path with unique name
            shutil.move(str(expected_docx_path), str(output_path))
            return output_path
        else:
            # Check if LibreOffice put it elsewhere
            fallback_docx = Path(f"{doc_path.stem}.docx")
            if fallback_docx.exists() and fallback_docx.stat().st_size > 0:
                shutil.move(str(fallback_docx), str(output_path))
                return output_path
            
            raise Exception(f"LibreOffice conversion failed - output DOCX not found at expected location")
            
    except subprocess.TimeoutExpired:
        raise Exception("LibreOffice conversion timed out after 120 seconds")
    except FileNotFoundError:
        raise Exception("LibreOffice is not installed or not found in PATH")
    except Exception as e:
        raise Exception(f"LibreOffice conversion failed: {str(e)}")


def convert_docx_to_pdf(docx_path: Path) -> Path:
    """Convert DOCX to PDF using LibreOffice (preferred), docx2pdf, or reportlab fallback.
    
    Conversion priority:
    1. LibreOffice - Best for Linux, native support without Wine
    2. docx2pdf - Works on Windows with MS Word or Linux with Wine
    3. reportlab - Basic text-only fallback
    """
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pdf"
    
    # Method 1: Try LibreOffice first (best for Linux)
    try:
        print("Attempting DOCX to PDF conversion using LibreOffice...")
        return convert_docx_to_pdf_libreoffice(docx_path)
    
    except Exception as e:
        print(f"LibreOffice conversion failed: {e}")
    
    # Method 2: Try docx2pdf
    try:
        print("Attempting DOCX to PDF conversion using docx2pdf...")
        docx2pdf_convert(str(docx_path), str(output_path))
        
        # Verify the output was created
        if output_path.exists() and output_path.stat().st_size > 0:
            print("docx2pdf conversion successful")
            return output_path
        else:
            raise Exception("docx2pdf conversion failed - output file not created or empty")
            
    except Exception as e:
        print(f"docx2pdf conversion failed: {e}")
    
    # Method 3: Fallback to reportlab (basic text-only conversion)
    print("Falling back to reportlab for basic text-only PDF conversion...")
    
    doc = Document(docx_path)
    pdf_canvas = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    y_position = height - 50
    
    for para in doc.paragraphs:
        if para.text.strip():
            text = para.text
            pdf_canvas.drawString(50, y_position, text[:100])
            y_position -= 20
            if y_position < 50:
                pdf_canvas.showPage()
                y_position = height - 50
    
    pdf_canvas.save()
    return output_path


def convert_doc_to_pdf(doc_path: Path) -> Path:
    """Convert DOC to PDF (via DOCX)"""
    docx_path = convert_doc_to_docx(doc_path)
    output_path = convert_docx_to_pdf(docx_path)
    return output_path


# ============== PDF Conversion Functions ==============

def convert_pdf_to_docx(pdf_path: Path) -> Path:
    """Convert PDF to DOCX using pdf2docx for better format preservation"""
    output_path = TEMP_DIR / f"{uuid.uuid4()}.docx"

    try:
        # Use pdf2docx Converter for better format preservation
        cv = Converter(str(pdf_path))
        cv.convert(str(output_path), start=0, end=None)
        cv.close()
    except Exception as e:
        # Fallback to text extraction if pdf2docx fails
        print(f"pdf2docx conversion failed: {e}, falling back to text extraction")
        reader = PdfReader(str(pdf_path))
        doc = Document()

        for page in reader.pages:
            text = page.extract_text()
            if text:
                lines = text.split("\n")
                for line in lines:
                    if line.strip():
                        doc.add_paragraph(line.strip())

        doc.save(output_path)

    return output_path


def convert_pdf_to_doc(pdf_path: Path) -> Path:
    """Convert PDF to DOC (via DOCX)"""
    docx_path = convert_pdf_to_docx(pdf_path)
    output_path = convert_docx_to_doc(docx_path)
    return output_path


def convert_pdf_to_text(pdf_path: Path) -> Path:
    """Extract text from PDF"""
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}.txt"
    output_path.write_text(text)
    return output_path


def convert_pdf_to_excel(pdf_path: Path, quality: str = "precise") -> Path:
    """Convert PDF to Excel using TableExtractionService for precise table extraction.
    
    Uses advanced table extraction with Camelot, pdfplumber, and Tabula for best results.
    
    Args:
        pdf_path: Path to the PDF file
        quality: Extraction quality - "precise" (uses Camelot, best accuracy) or "fast" (uses pdfplumber, faster)
    
    Returns:
        Path to the converted Excel file
    """
    # Define border style for table cells
    thin_border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )

    wb = Workbook()
    
    # Initialize the table extraction service
    # prefer_quality=True means accuracy over speed
    extraction_service = TableExtractionService(prefer_quality=(quality == "precise"))
    
    # Extract all tables using the advanced service
    print(f"Extracting tables from PDF with {quality} quality...")
    all_tables = extraction_service.extract_all_tables(pdf_path)
    
    # Also extract text content for non-table areas
    with pdfplumber.open(pdf_path) as pdf:
        total_pages = len(pdf.pages)
        
        for page_num, page in enumerate(pdf.pages, start=1):
            if page_num == 1:
                ws = wb.active
                ws.title = f"Page {page_num}"
            else:
                ws = wb.create_sheet(f"Page {page_num}")
            
            current_row = 1
            
            # Extract and embed images first
            try:
                images = page.images
                if images:
                    for img_idx, img in enumerate(images):
                        try:
                            # Extract image data
                            img_bbox = (img['x0'], img['top'], img['x1'], img['bottom'])
                            img_data = page.within_bbox(img_bbox).to_image(resolution=150)
                            
                            # Save image to temporary file
                            img_temp_path = TEMP_DIR / f"temp_img_{page_num}_{img_idx}.png"
                            img_data.save(img_temp_path)
                            
                            # Embed image in Excel
                            from openpyxl.drawing.image import Image as XLImage
                            xl_img = XLImage(img_temp_path)
                            
                            # Resize image to fit reasonably in Excel (max 200x200 pixels)
                            xl_img.width = min(xl_img.width, 200)
                            xl_img.height = min(xl_img.height, 200)
                            
                            # Add image to worksheet
                            ws.add_image(xl_img, f"A{current_row}")
                            
                            # Add label for the image
                            ws.cell(row=current_row, column=2, value=f"Image {img_idx + 1} from Page {page_num}")
                            ws.cell(row=current_row, column=2).font = Font(italic=True)
                            
                            current_row += max(15, int(xl_img.height / 15) + 2)
                            
                        except Exception as e:
                            print(f"Failed to extract image {img_idx} from page {page_num}: {e}")
                            continue
            except Exception as e:
                print(f"Failed to extract images from page {page_num}: {e}")
            
            # Process tables extracted by TableExtractionService
            if page_num in all_tables and all_tables[page_num]:
                tables_on_page = all_tables[page_num]
                
                for table_idx, table_data in enumerate(tables_on_page):
                    if table_idx > 0 or current_row > 1:
                        ws.cell(row=current_row, column=1, value="")
                        current_row += 1
                    
                    # Add table metadata header
                    metadata = table_data.metadata
                    if metadata.table_title:
                        ws.cell(row=current_row, column=1, value=metadata.table_title)
                        ws.cell(row=current_row, column=1).font = Font(bold=True, size=12)
                        current_row += 1
                    
                    # Add confidence score info
                    confidence_pct = metadata.confidence_score * 100
                    ws.cell(row=current_row, column=1, value=f"Confidence: {confidence_pct:.1f}% | Method: {metadata.method_used.value}")
                    ws.cell(row=current_row, column=1).font = Font(italic=True, size=10)
                    current_row += 1
                    
                    # Add extraction warnings if any
                    if metadata.extraction_warnings:
                        ws.cell(row=current_row, column=1, value=f"Warnings: {', '.join(metadata.extraction_warnings)}")
                        ws.cell(row=current_row, column=1).font = Font(italic=True, size=9, color="FFA500")
                        current_row += 1
                    
                    # Get the table data
                    table_cells = table_data.cells
                    
                    if table_cells:
                        table_start_row = current_row
                        
                        for row_idx, row in enumerate(table_cells):
                            for col_idx, cell in enumerate(row, start=1):
                                excel_cell = ws.cell(row=current_row, column=col_idx, value=cell.formatted_value)
                                
                                # Apply border
                                excel_cell.border = thin_border
                                
                                # Center align
                                excel_cell.alignment = Alignment(horizontal='center', vertical='center')
                                
                                # Make header row bold
                                if cell.is_header or (row_idx == 0 and metadata.has_header):
                                    excel_cell.font = Font(bold=True)
                                
                                # Apply data type specific formatting
                                if cell.data_type == "currency":
                                    excel_cell.number_format = '"$"#,##0.00'
                                elif cell.data_type == "percentage":
                                    excel_cell.number_format = '0.00%'
                                elif cell.data_type == "number":
                                    excel_cell.number_format = '#,##0.00'
                            
                            current_row += 1
                        
                        # Adjust column widths based on content
                        for col in range(1, len(table_cells[0]) + 1):
                            max_width = 0
                            for row in range(table_start_row, current_row):
                                cell = ws.cell(row=row, column=col)
                                if cell.value:
                                    cell_width = len(str(cell.value)) + 2
                                    if cell_width > max_width:
                                        max_width = cell_width
                            if max_width > 50:
                                max_width = 50
                            if max_width < 8:
                                max_width = 8
                            col_letter = ws.cell(row=table_start_row, column=col).column_letter
                            ws.column_dimensions[col_letter].width = max_width
            
            # Fallback: use basic pdfplumber extraction if no tables found
            if page_num not in all_tables or not all_tables[page_num]:
                tables = page.extract_tables()
                
                if tables:
                    if current_row > 1:
                        ws.cell(row=current_row, column=1, value="")
                        current_row += 1
                    
                    for table_idx, table in enumerate(tables):
                        if table_idx > 0 or current_row > 1:
                            ws.cell(row=current_row, column=1, value="")
                            current_row += 1
                        
                        ws.cell(row=current_row, column=1, value=f"Table {table_idx + 1} (basic extraction)")
                        ws.cell(row=current_row, column=1).font = Font(bold=True, size=12)
                        current_row += 1
                        
                        table_start_row = current_row
                        
                        for row_idx, row in enumerate(table):
                            if row is None:
                                continue
                            
                            for col_idx, cell_value in enumerate(row, start=1):
                                cell = ws.cell(row=current_row, column=col_idx, value=cell_value.strip() if cell_value else "")
                                cell.border = thin_border
                                cell.alignment = Alignment(horizontal='center', vertical='center')
                                if row_idx == 0:
                                    cell.font = Font(bold=True)
                            
                            current_row += 1
                        
                        for col in range(1, ws.max_column + 1):
                            max_width = 0
                            for row in range(table_start_row, current_row):
                                cell = ws.cell(row=row, column=col)
                                if cell.value:
                                    cell_width = len(str(cell.value)) + 2
                                    if cell_width > max_width:
                                        max_width = cell_width
                            if max_width > 50:
                                max_width = 50
                            col_letter = ws.cell(row=table_start_row, column=col).column_letter
                            ws.column_dimensions[col_letter].width = max_width
            
            # Extract text content for non-table areas
            text = page.extract_text()
            if text and current_row > 1:
                ws.cell(row=current_row, column=1, value="")
                current_row += 1
            
            if text:
                ws.cell(row=current_row, column=1, value=f"Page {page_num} - Text Content:")
                ws.cell(row=current_row, column=1).font = Font(bold=True, size=12)
                current_row += 1
                
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        ws.cell(row=current_row, column=1, value=line.strip())
                        current_row += 1
            
            # If page has no content at all
            if current_row == 1:
                ws.cell(row=1, column=1, value=f"Page {page_num} - No extractable content found")
                ws.cell(row=1, column=1).font = Font(italic=True)
            
            print(f"Processed page {page_num}/{total_pages}")
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}.xlsx"
    wb.save(output_path)
    print(f"Excel file saved to: {output_path}")
    return output_path


# ============== PDF to PPTX Conversion Helper Functions ==============

    wb = Workbook()

    # Open PDF with pdfplumber for comprehensive extraction
    with pdfplumber.open(pdf_path) as pdf:
        # Create a sheet for each page
        for page_num, page in enumerate(pdf.pages, start=1):
            if page_num == 1:
                ws = wb.active
                ws.title = f"Page {page_num}"
            else:
                ws = wb.create_sheet(f"Page {page_num}")

            current_row = 1

            # Extract and embed images first
            try:
                images = page.images
                if images:
                    for img_idx, img in enumerate(images):
                        try:
                            # Extract image data
                            img_bbox = (img['x0'], img['top'], img['x1'], img['bottom'])
                            img_data = page.within_bbox(img_bbox).to_image(resolution=150)

                            # Save image to temporary file
                            img_temp_path = TEMP_DIR / f"temp_img_{page_num}_{img_idx}.png"
                            img_data.save(img_temp_path)

                            # Embed image in Excel
                            from openpyxl.drawing.image import Image as XLImage
                            xl_img = XLImage(img_temp_path)

                            # Resize image to fit reasonably in Excel (max 200x200 pixels)
                            xl_img.width = min(xl_img.width, 200)
                            xl_img.height = min(xl_img.height, 200)

                            # Add image to worksheet
                            ws.add_image(xl_img, f"A{current_row}")

                            # Add label for the image
                            ws.cell(row=current_row, column=2, value=f"Image {img_idx + 1} from Page {page_num}")
                            ws.cell(row=current_row, column=2).font = Font(italic=True)

                            current_row += max(15, int(xl_img.height / 15) + 2)  # Space for image + label

                        except Exception as e:
                            print(f"Failed to extract image {img_idx} from page {page_num}: {e}")
                            continue
            except Exception as e:
                print(f"Failed to extract images from page {page_num}: {e}")

            # Extract tables from the page
            tables = page.extract_tables()

            if tables:
                # Process each table found on the page
                for table_idx, table in enumerate(tables):
                    if table_idx > 0 or current_row > 1:
                        # Add a blank row between tables/content
                        ws.cell(row=current_row, column=1, value="")
                        current_row += 1

                    # Add table header
                    ws.cell(row=current_row, column=1, value=f"Table {table_idx + 1}")
                    ws.cell(row=current_row, column=1).font = Font(bold=True, size=12)
                    current_row += 1

                    table_start_row = current_row

                    for row_idx, row in enumerate(table):
                        if row is None:
                            continue

                        for col_idx, cell_value in enumerate(row, start=1):
                            cell = ws.cell(row=current_row, column=col_idx, value=cell_value.strip() if cell_value else "")

                            # Apply border to all cells
                            cell.border = thin_border

                            # Center align cells
                            cell.alignment = Alignment(horizontal='center', vertical='center')

                            # Make header row bold (first row of each table)
                            if row_idx == 0:
                                cell.font = Font(bold=True)

                        current_row += 1

                    # Adjust column widths based on content for this table
                    for col in range(1, ws.max_column + 1):
                        max_width = 0
                        for row in range(table_start_row, current_row):
                            cell = ws.cell(row=row, column=col)
                            if cell.value:
                                cell_width = len(str(cell.value)) + 2
                                if cell_width > max_width:
                                    max_width = cell_width
                        # Limit column width to reasonable values
                        if max_width > 50:
                            max_width = 50
                        col_letter = ws.cell(row=table_start_row, column=col).column_letter
                        current_width = ws.column_dimensions.get(col_letter, None)
                        current_width = current_width.width if current_width else 0
                        ws.column_dimensions[col_letter].width = max(current_width, min(max_width, 50))

            # Extract text content if no tables found or as additional content
            text = page.extract_text()
            if text and (not tables or current_row > 1):
                # Add a blank row before text section
                if current_row > 1:
                    ws.cell(row=current_row, column=1, value="")
                    current_row += 1

                # Add text header
                ws.cell(row=current_row, column=1, value=f"Page {page_num} - Text Content:")
                ws.cell(row=current_row, column=1).font = Font(bold=True, size=12)
                current_row += 1

                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        ws.cell(row=current_row, column=1, value=line.strip())
                        current_row += 1

            # If page has no content at all
            if current_row == 1:
                ws.cell(row=1, column=1, value=f"Page {page_num} - No extractable content found")
                ws.cell(row=1, column=1).font = Font(italic=True)



def ocr_page(page):
    image = page.to_image(resolution=300).original
    return pytesseract.image_to_string(image)

def detect_columns(words, tolerance=30):
    columns = {}
    for w in words:
        x = round(w["x0"] / tolerance) * tolerance
        columns.setdefault(x, []).append(w)
    return list(columns.values())

def overlay_text(slide, words):
    lines = {}
    for w in words:
        y = round(w["top"], 1)
        lines.setdefault(y, []).append(w)

    for y, line_words in lines.items():
        line_words.sort(key=lambda x: x["x0"])
        text = " ".join(w["text"] for w in line_words)

        left = Inches(min(w["x0"] for w in line_words) / 72)
        top = Inches(y / 72)
        width = Inches(
            (max(w["x1"] for w in line_words) -
             min(w["x0"] for w in line_words)) / 72
        )

        box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.name = "Calibri"

def add_background_image(slide, pdf_path, page_num, slide_width, slide_height):
    images = convert_from_path(
        str(pdf_path),
        dpi=300,
        first_page=page_num,
        last_page=page_num
    )
    img_path = TEMP_DIR / f"bg_{uuid.uuid4()}.png"
    images[0].save(img_path)

    slide.shapes.add_picture(
        str(img_path), Inches(0), Inches(0),
        width=slide_width, height=slide_height
    )
    img_path.unlink(missing_ok=True)

def compute_quality_score(used_ocr, detected_tables, detected_columns, word_count):
    score = 100
    if used_ocr:
        score -= 30
    if detected_tables:
        score += 10
    if not detected_columns:
        score -= 10
    if word_count < 20:
        score -= 20
    return max(0, min(score, 100))


# ===================== MAIN CONVERTER  PDF TO PPTX=====================
from pptx.util import Inches
from pathlib import Path
import uuid
import numpy as np
import pdfplumber
import pytesseract
import camelot
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pdf2image import convert_from_path
from pypdf import PdfReader

# def convert_pdf_to_pptx(pdf_path: Path) -> Path:
#     """
#     Adobe-like PDF → PPTX converter:
#     - Two-layer rendering
#     - OCR fallback
#     - Column detection
#     - Table detection
#     - Quality scoring
#     """

#     reader = PdfReader(str(pdf_path))
#     first_page = reader.pages[0]

#     pdf_w = float(first_page.mediabox.width) / 72
#     pdf_h = float(first_page.mediabox.height) / 72

#     prs = Presentation()
#     prs.slide_width = Inches(pdf_w)
#     prs.slide_height = Inches(pdf_h)

#     quality_scores = []

#     with pdfplumber.open(pdf_path) as pdf:
#         for idx, page in enumerate(pdf.pages, start=1):
#             slide = prs.slides.add_slide(prs.slide_layouts[6])

#             # ---------- Layer 0 (Background Image) ----------
#             add_background_image(
#                 slide, pdf_path, idx,
#                 prs.slide_width, prs.slide_height
#             )

#             words = page.extract_words(use_text_flow=True)
#             used_ocr = False

#             if not words:
#                 text = ocr_page(page)
#                 used_ocr = True
#                 words = [
#                     {"text": t, "x0": 50, "x1": 500, "top": i * 14}
#                     for i, t in enumerate(text.splitlines())
#                     if t.strip()
#                 ]

#             # ---------- Column Detection ----------
#             columns = detect_columns(words)

#             # ---------- Editable Overlay ----------
#             overlay_text(slide, words)

#             # ---------- Table Detection ----------
#             tables = []
#             try:
#                 tables = camelot.read_pdf(
#                     str(pdf_path),
#                     pages=str(idx),
#                     flavor="stream"
#                 )
#             except Exception:
#                 pass

#             for table in tables:
#                 rows, cols = table.df.shape
#                 ppt_table = slide.shapes.add_table(
#                     rows, cols,
#                     Inches(0.5), Inches(0.5),
#                     Inches(6), Inches(3)
#                 ).table

#                 for r in range(rows):
#                     for c in range(cols):
#                         ppt_table.cell(r, c).text = table.df.iloc[r, c]

#             # ---------- Quality Score ----------
#             score = compute_quality_score(
#                 used_ocr=used_ocr,
#                 detected_tables=bool(tables),
#                 detected_columns=len(columns) > 1,
#                 word_count=len(words)
#             )
#             quality_scores.append(score)

#     output_path = TEMP_DIR / f"{uuid.uuid4()}.pptx"
#     prs.save(output_path)

#     print("Average quality score:", sum(quality_scores) // len(quality_scores))
#     return output_path





def convert_pdf_to_pptx(pdf_path: Path) -> Path:
    """Convert PDF to PowerPoint with enhanced layout and media preservation.
    
    This function provides two conversion methods:
    1. Hybrid-based (enhanced): Extracts text with layout AND embedded media
       - Preserves visual elements (images, charts, formatting, layout)
       - Editable text with proper formatting
       - Best balance of visual fidelity and editability
    2. Image-based: Converts each PDF page to an image and adds as slide
       - Preserves all visual elements exactly as shown
       - Best for visual fidelity when editability is not needed
    3. Text-based: Extracts text with basic layout information
       - Editable text output
       - Works without external dependencies
       - Good for text-heavy documents
    
    Returns:
        Path: Path to the converted PowerPoint file
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    
    # Try hybrid-based conversion first (best balance of quality and editability)
    try:
        print("Attempting image-based PDF to PPTX conversion...")
        return _convert_pdf_to_pptx_image_based(pdf_path)

    except Exception as image_error:
        print(f"Image-based conversion failed: {image_error}, trying hybrid conversion...")

        try:
            return _convert_pdf_to_pptx_hybrid(pdf_path)

        except ImportError as e:
            print(f"pdfplumber not available: {e}, falling back to text-based conversion")
            return _convert_pdf_to_pptx_text_based(pdf_path)

        except Exception as hybrid_error:
            print(f"Hybrid conversion failed: {hybrid_error}, falling back to text-based conversion")
            return _convert_pdf_to_pptx_text_based(pdf_path)


def _convert_pdf_to_pptx_hybrid(pdf_path: Path) -> Path:
    """Convert PDF to PowerPoint using hybrid approach.
    
    This method extracts text with layout preservation AND preserves embedded media.
    It provides the best balance of visual fidelity and editability.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    # Open PDF with pdfplumber for comprehensive extraction
    with pdfplumber.open(pdf_path) as pdf:
        prs = Presentation()
        
        # Get source PDF page dimensions and calculate appropriate slide size
        first_page = pdf.pages[0]
        pdf_width_pt = first_page.width
        pdf_height_pt = first_page.height
        
        # Convert PDF points to inches (72 points per inch)
        pdf_width_inches = pdf_width_pt / 72
        pdf_height_inches = pdf_height_pt / 72
        
        # Set slide dimensions to match source PDF (maintain aspect ratio)
        # Use 16:9 as baseline but adjust based on source
        prs.slide_width = Inches(pdf_width_inches)
        prs.slide_height = Inches(pdf_height_inches)
        
        # Process each page
        for page_num, page in enumerate(pdf.pages, start=1):
            print(f"Processing page {page_num}/{len(pdf.pages)}...")
            
            # Create a blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            current_y = Inches(0.5)  # Start with margin from top
            left_margin = Inches(0.5)
            max_width = prs.slide_width - Inches(1)  # Left and right margins
            
            # 1. First, extract and add images from the page
            try:
                if page.images:
                    print(f"  Found {len(page.images)} images on page {page_num}")
                    for img_idx, img in enumerate(page.images):
                        try:
                            # Extract image data
                            img_bbox = (img['x0'], img['top'], img['x1'], img['bottom'])
                            img_data = page.within_bbox(img_bbox).to_image(resolution=300)
                            
                            # Save image to temporary file
                            img_temp_path = TEMP_DIR / f"temp_pptx_img_{page_num}_{img_idx}.png"
                            img_data.save(str(img_temp_path))
                            
                            # Calculate image dimensions
                            img_width_inches = (img['x1'] - img['x0']) / 72
                            img_height_inches = (img['bottom'] - img['top']) / 72
                            
                            # Scale down if image is wider than slide
                            if img_width_inches > max_width:
                                scale = max_width / Inches(img_width_inches)
                                img_width_inches = max_width
                                img_height_inches = Inches(float(img_height_inches) * float(scale))
                            
                            # Add image to slide
                            img_left = left_margin
                            img_top = current_y
                            
                            slide.shapes.add_picture(
                                str(img_temp_path), 
                                img_left, 
                                img_top, 
                                width=img_width_inches
                            )
                            
                            # Update Y position for next element
                            current_y = img_top + img_height_inches + Inches(0.3)
                            
                            # Clean up temp image
                            img_temp_path.unlink(missing_ok=True)
                            
                        except Exception as e:
                            print(f"  Failed to extract image {img_idx}: {e}")
                            continue
            except Exception as e:
                print(f"  Failed to extract images from page {page_num}: {e}")
            
            # 2. Extract and add text content
            text = page.extract_text()
            if text:
                # Clean up text
                lines = text.split('\n')
                cleaned_lines = []
                for line in lines:
                    line = line.strip()
                    if line:
                        cleaned_lines.append(line)
                
                if cleaned_lines:
                    # Create text box for content
                    text_box_height = prs.slide_height - current_y - Inches(0.5)
                    
                    if text_box_height > Inches(0.5):
                        text_box = slide.shapes.add_textbox(
                            left_margin,
                            current_y,
                            max_width,
                            text_box_height
                        )
                        tf = text_box.text_frame
                        tf.word_wrap = True
                        
                        # Process lines and add to text frame
                        for i, line in enumerate(cleaned_lines):
                            # Detect headers (lines that are short and followed by longer content)
                            is_header = (
                                len(line) < 50 and 
                                i < len(cleaned_lines) - 1 and 
                                len(cleaned_lines[i + 1]) > len(line)
                            )
                            
                            # Detect bullet points
                            is_bullet = (
                                line.startswith('•') or 
                                line.startswith('·') or
                                line.startswith('- ') or
                                line.startswith('* ') or
                                (len(line) > 2 and line[0].isdigit() and line[1] in '.)')
                            )
                            
                            # Clean bullet prefix
                            if is_bullet:
                                for prefix in ['• ', '· ', '- ', '* ', '1. ', '2. ', '3. ', '1) ', '2) ', '3) ']:
                                    if line.startswith(prefix):
                                        line = line[len(prefix):]
                                        break
                            
                            if i == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            
                            p.text = line
                            p.font.size = Pt(12)
                            
                            # Apply formatting based on content type
                            if is_header:
                                p.font.size = Pt(18)
                                p.font.bold = True
                                p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
                                p.space_before = Pt(12)
                                p.space_after = Pt(6)
                            elif is_bullet:
                                p.level = 0
                                p.space_before = Pt(6)
                            else:
                                p.space_before = Pt(4)
                        
                        print(f"  Added text content with {len(cleaned_lines)} lines")
            
            print(f"  Completed page {page_num}")
        
        output_path = TEMP_DIR / f"{uuid.uuid4()}.pptx"
        prs.save(output_path)
        print(f"PPTX saved to: {output_path}")
        return output_path


def _convert_pdf_to_pptx_image_based(pdf_path: Path) -> Path:
    """Convert PDF to PowerPoint using image-based approach.
    
    Each PDF page is converted to a high-quality image and added as a slide.
    This preserves visual layout, images, and formatting perfectly.
    Uses 300 DPI for high-quality output.
    
    Key improvements:
    - Slide dimensions match the source PDF page dimensions exactly
    - Images are scaled to fit while maintaining aspect ratio
    - High-quality image conversion with proper temp file cleanup
    """
    from pdf2image import convert_from_path
    from pptx.util import Inches
    from PIL import Image as PILImage
    
    # Get PDF page dimensions first
    from pypdf import PdfReader
    reader = PdfReader(str(pdf_path))
    first_page = reader.pages[0]
    
    # Get PDF page dimensions in points (72 points per inch)
    pdf_width_pt = float(first_page.mediabox.width)
    pdf_height_pt = float(first_page.mediabox.height)
    
    # Convert to inches for PPTX
    pdf_width_inches = pdf_width_pt / 72
    pdf_height_inches = pdf_height_pt / 72
    
    print(f"PDF page dimensions: {pdf_width_pt}pt x {pdf_height_pt}pt ({pdf_width_inches:.2f}\" x {pdf_height_inches:.2f}\")")
    
    # Convert PDF pages to images at 300 DPI for high quality
    print(f"Converting PDF to images (300 DPI): {pdf_path}")
    images = convert_from_path(str(pdf_path), dpi=300, thread_count=2)
    
    prs = Presentation()
    
    # Set slide dimensions to match the source PDF page dimensions exactly
    # This preserves the aspect ratio of the original document
    prs.slide_width = Inches(pdf_width_inches)
    prs.slide_height = Inches(pdf_height_inches)
    
    print(f"Slide dimensions set to: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    
    temp_images = []  # Track temp files for cleanup
    
    for page_num, image in enumerate(images):
        try:
            # Create blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Get image dimensions
            img_width_px, img_height_px = image.size
            
            # Calculate scaling to fit image on slide while maintaining aspect ratio
            # The image should fill the entire slide
            slide_width_px = int(prs.slide_width.inches * 300)  # 300 DPI
            slide_height_px = int(prs.slide_height.inches * 300)
            
            # Calculate the scale factors
            scale_w = slide_width_px / img_width_px
            scale_h = slide_height_px / img_height_px
            
            # Use the larger scale to fill the slide completely (cover mode)
            # This ensures no white space around the image
            scale = max(scale_w, scale_h)
            
            # Calculate scaled dimensions
            scaled_width = int(img_width_px * scale)
            scaled_height = int(img_height_px * scale)
            
            # Calculate position to center the image
            left_px = (slide_width_px - scaled_width) // 2
            top_px = (slide_height_px - scaled_height) // 2
            
            # Convert pixels back to inches for PPTX
            left = Inches(left_px / 300)
            top = Inches(top_px / 300)
            width = Inches(scaled_width / 300)
            
            # Save image temporarily at high quality
            img_path = TEMP_DIR / f"temp_pptx_img_{uuid.uuid4()}.png"
            temp_images.append(img_path)
            image.save(str(img_path), "PNG", quality=95)
            
            # Add image to slide
            slide.shapes.add_picture(str(img_path), left, top, width=width)
            
            print(f"Added slide {page_num + 1}/{len(images)}: {img_width_px}x{img_height_px}px -> {scaled_width}x{scaled_height}px")
            
        except Exception as e:
            print(f"Error processing page {page_num + 1}: {e}")
            continue
    
    # Clean up temp images
    for img_path in temp_images:
        try:
            if img_path.exists():
                img_path.unlink()
        except Exception as e:
            print(f"Warning: Failed to clean up temp file {img_path}: {e}")
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pptx"
    prs.save(output_path)
    print(f"PPTX saved to: {output_path}")
    return output_path


def _convert_pdf_to_pptx_text_based(pdf_path: Path) -> Path:
    """Convert PDF to PowerPoint using text extraction with layout preservation.
    
    Extracts text from PDF while preserving some layout information.
    Creates editable slides with extracted text.
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    reader = PdfReader(str(pdf_path))
    prs = Presentation()
    
    # Use widescreen (16:9) format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for page_num, page in enumerate(reader.pages):
        # Create a slide with title layout
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Get page dimensions
        page_width_pt = float(page.mediabox.width)
        page_height_pt = float(page.mediabox.height)
        
        # Extract text from page
        text = page.extract_text()
        
        # Add page title
        title = slide.shapes.title
        title.text = f"Page {page_num + 1}"
        if title.text_frame:
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.bold = True
        
        # Get content placeholder
        content_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        if content_shape and text:
            text_frame = content_shape.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            # Split text into paragraphs and add to slide
            paragraphs = text.split('\n')
            
            for i, para in enumerate(paragraphs):
                para = para.strip()
                if not para:
                    continue
                
                # Truncate very long paragraphs
                if len(para) > 500:
                    para = para[:497] + "..."
                
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = para
                p.font.size = Pt(14)
                p.space_before = Pt(6)
                
                # Add bullet points for shorter items (likely list items)
                if len(para) < 100 and (para.startswith("•") or para.startswith("-") or para.startswith("·") or para[0:2].replace(".", "").isdigit()):
                    p.level = 0
        elif text:
            # If no content placeholder, create a text box
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12.333)
            height = Inches(5.5)
            
            text_box = slide.shapes.add_textbox(left, top, width, height)
            tf = text_box.text_frame
            tf.word_wrap = True
            
            paragraphs = text.split('\n')
            for i, para in enumerate(paragraphs):
                para = para.strip()
                if not para:
                    continue
                
                if len(para) > 500:
                    para = para[:497] + "..."
                
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = para
                p.font.size = Pt(14)
                p.space_before = Pt(6)
        
        print(f"Processed page {page_num + 1}/{len(reader.pages)}")
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pptx"
    prs.save(output_path)
    print(f"PPTX saved to: {output_path}")
    return output_path


# ============== Excel Conversion Functions ==============

def convert_excel_to_pdf_libreoffice(excel_path: Path) -> Path:
    """Convert Excel to PDF using LibreOffice (best format preservation).
    
    This method provides the best quality PDF output with:
    - Perfect table border preservation
    - Cell colors and backgrounds maintained
    - Column widths and row heights preserved
    - Font formatting (bold, italic, colors)
    - Merged cells and complex layouts
    - Multiple sheets handled correctly
    
    Uses LibreOffice in headless mode which provides native Excel rendering.
    """
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pdf"
    
    try:
        # LibreOffice requires a HOME directory to work properly
        env = os.environ.copy()
        if 'HOME' not in env:
            env['HOME'] = '/tmp'
        
        # LibreOffice command: headless mode, convert to PDF, specify output directory
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(TEMP_DIR),
            str(excel_path)
        ]
        
        print(f"Running LibreOffice command: {' '.join(cmd)}")
        
        # Run the conversion
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120,  # 2 minute timeout
            env=env
        )
        
        if result.returncode != 0:
            print(f"LibreOffice stderr: {result.stderr}")
            raise Exception(f"LibreOffice conversion failed with return code {result.returncode}")
        
        # LibreOffice creates the PDF with the same name as input but .pdf extension
        # in the output directory
        expected_pdf_path = TEMP_DIR / f"{excel_path.stem}.pdf"
        
        if expected_pdf_path.exists() and expected_pdf_path.stat().st_size > 0:
            # Move to our output path with unique name
            shutil.move(str(expected_pdf_path), str(output_path))
            print(f"LibreOffice conversion successful: {output_path}")
            return output_path
        else:
            # Check if LibreOffice put it elsewhere
            fallback_pdf = Path(f"{excel_path.stem}.pdf")
            if fallback_pdf.exists() and fallback_pdf.stat().st_size > 0:
                shutil.move(str(fallback_pdf), str(output_path))
                print(f"LibreOffice conversion successful (fallback location): {output_path}")
                return output_path
            
            raise Exception(f"LibreOffice conversion failed - output PDF not found at expected location")
            
    except subprocess.TimeoutExpired:
        raise Exception("LibreOffice conversion timed out after 120 seconds")
    except FileNotFoundError:
        raise Exception("LibreOffice is not installed or not found in PATH")
    except Exception as e:
        raise Exception(f"LibreOffice conversion failed: {str(e)}")


def check_libreoffice_available() -> bool:
    """Check if LibreOffice is available on the system."""
    try:
        result = subprocess.run(
            ['libreoffice', '--version'],
            capture_output=True,
            text=True,
            timeout=10
        )
        if result.returncode == 0:
            print(f"LibreOffice available: {result.stdout.strip()}")
            return True
    except (subprocess.SubprocessError, FileNotFoundError):
        pass
    return False


def _xlsx2pdf_conversion(excel_path: Path, output_path: Path) -> Path:
    """Fallback conversion using xlsx2pdf when LibreOffice is not available."""
    try:
        # Find a suitable font for the Transformer
        font_path = find_font_path()
        
        # Initialize the Transformer with the font path
        transformer = Transformer(font_path)
        
        # Load the workbook using openpyxl
        wb = load_workbook(excel_path)
        ws = wb.active
        
        # Get the dimensions of the data
        rows_n = ws.max_row
        cols_n = ws.max_column
        
        print(f"Converting Excel file: {excel_path}")
        print(f"Dimensions: {rows_n} rows x {cols_n} columns")
        
        # Call transform() which returns PDF bytes
        pdf_bytes = transformer.transform(wb, rows_n, cols_n)
        
        # Write the PDF bytes to the output file
        with open(output_path, "wb") as f:
            f.write(pdf_bytes)
        
        # Verify the output was created
        if output_path.exists() and output_path.stat().st_size > 0:
            print(f"Excel to PDF conversion successful using xlsx2pdf: {output_path}")
            return output_path
        else:
            raise Exception("xlsx2pdf conversion failed - output file not created or empty")
            
    except Exception as e:
        raise Exception(f"xlsx2pdf conversion failed: {str(e)}")


def convert_excel_to_pdf(excel_path: Path) -> Path:
    """Convert Excel file to PDF with best format preservation.
    
    Conversion priority (best quality first):
    1. LibreOffice - Best for format preservation (primary method)
       - Preserves ALL Excel formatting
       - Native Excel rendering through LibreOffice
       - Perfect table borders, colors, fonts, column widths
    
    2. xlsx2pdf - Secondary fallback
       - Basic formatting support
       - Good for simple Excel files
    
    3. reportlab - Last resort fallback
       - Basic text-only conversion
       - No formatting preserved
    
    Returns:
        Path: Path to the converted PDF file
    """
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pdf"
    
    # Method 1: Try LibreOffice first (best format preservation)
    try:
        print("Attempting Excel to PDF conversion using LibreOffice...")
        return convert_excel_to_pdf_libreoffice(excel_path)
    except Exception as libreoffice_error:
        print(f"LibreOffice conversion failed: {libreoffice_error}")
    
    # Method 2: Try xlsx2pdf as fallback
    try:
        print("Attempting Excel to PDF conversion using xlsx2pdf...")
        return _xlsx2pdf_conversion(excel_path, output_path)
    except Exception as xlsx2pdf_error:
        print(f"xlsx2pdf conversion failed: {xlsx2pdf_error}")
    
    # Method 3: Basic reportlab fallback (last resort)
    print("Falling back to basic reportlab conversion...")
    print("WARNING: This method does not preserve formatting!")
    return _fallback_excel_to_pdf(excel_path, output_path)


def _fallback_excel_to_pdf(excel_path: Path, output_path: Path) -> Path:
    """Fallback conversion using reportlab when xlsx2pdf fails."""
    wb = load_workbook(excel_path)
    ws = wb.active
    
    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    y_position = height - 50
    line_height = 12
    
    max_row = ws.max_row
    max_col = ws.max_column
    
    c.setFont("Helvetica-Bold", 14)
    c.drawString(50, y_position, f"Excel Document: {ws.title}")
    y_position -= 30
    
    c.setFont("Helvetica", 10)
    
    for row in range(1, max_row + 1):
        row_text = []
        for col in range(1, max_col + 1):
            cell = ws.cell(row=row, column=col)
            value = cell.value if cell.value is not None else ""
            row_text.append(str(value))
        
        text = " | ".join(row_text)
        
        if y_position < 50:
            c.showPage()
            y_position = height - 50
        
        if len(text) > 100:
            text = text[:97] + "..."
        
        c.drawString(50, y_position, text)
        y_position -= line_height
    
    c.save()
    print(f"Fallback conversion completed: {output_path}")
    return output_path


# ============== PowerPoint Conversion Functions ==============

def _check_libreoffice():
    """Ensure LibreOffice is installed"""
    if not shutil.which("libreoffice"):
        raise EnvironmentError(
            "LibreOffice not found. Install it using:\n"
            "sudo apt install libreoffice"
        )


def convert_pptx_to_pdf(pptx_path: Path) -> Path:
    """
    Convert PPTX to PDF while preserving layout, fonts, images, charts.
    Uses LibreOffice headless rendering engine.
    """
    _check_libreoffice()

    pptx_path = pptx_path.resolve()

    if not pptx_path.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")

    if pptx_path.suffix.lower() != ".pptx":
        raise ValueError("Expected a .pptx file")

    # Unique output folder to avoid name clashes
    output_dir = TEMP_DIR / str(uuid.uuid4())
    output_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        "libreoffice",
        "--headless",
        "--invisible",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        "--convert-to", "pdf",
        str(pptx_path),
        "--outdir", str(output_dir),
    ]

    subprocess.run(
        cmd,
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.PIPE,
        timeout=120
    )

    pdf_path = output_dir / f"{pptx_path.stem}.pdf"

    if not pdf_path.exists():
        raise RuntimeError("PDF conversion failed")

    return pdf_path


def convert_ppt_to_pdf(ppt_path: Path) -> Path:
    """
    Convert PPT to PDF.
    LibreOffice automatically handles .ppt and .pptx.
    """
    _check_libreoffice()

    ppt_path = ppt_path.resolve()

    if not ppt_path.exists():
        raise FileNotFoundError(f"File not found: {ppt_path}")

    if ppt_path.suffix.lower() not in {".ppt", ".pptx"}:
        raise ValueError("Only .ppt and .pptx files are supported")

    output_dir = TEMP_DIR / str(uuid.uuid4())
    output_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        "libreoffice",
        "--headless",
        "--invisible",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        "--convert-to", "pdf",
        str(ppt_path),
        "--outdir", str(output_dir),
    ]

    subprocess.run(
        cmd,
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.PIPE,
        timeout=120
    )

    pdf_path = output_dir / f"{ppt_path.stem}.pdf"

    if not pdf_path.exists():
        raise RuntimeError("PDF conversion failed")

    return pdf_path


# ============== Text Conversion Functions ==============

def convert_text_to_docx(text_path: Path) -> Path:
    """Convert text file to DOCX"""
    text = text_path.read_text()
    doc = Document()
    doc.add_paragraph(text)
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}.docx"
    doc.save(output_path)
    return output_path


def convert_text_to_pdf(text_path: Path) -> Path:
    """Convert text file to PDF"""
    text = text_path.read_text()
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pdf"
    
    pdf_canvas = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    y_position = height - 50
    
    for line in text.split('\n'):
        if line.strip():
            pdf_canvas.drawString(50, y_position, line[:100])
            y_position -= 15
            if y_position < 50:
                pdf_canvas.showPage()
                y_position = height - 50
    
    pdf_canvas.save()
    return output_path


# ============== Image Conversion Functions ==============

# Define supported image formats
SUPPORTED_IMAGE_FORMATS = {
    # Raster formats
    'jpg': 'JPEG',
    'jpeg': 'JPEG',
    'png': 'PNG',
    'gif': 'GIF',
    'bmp': 'BMP',
    'tiff': 'TIFF',
    'tif': 'TIFF',
    'webp': 'WEBP',
    'ico': 'ICO',
    'pcx': 'PCX',
    'ppm': 'PPM',
    'pgm': 'PGM',
    'pbm': 'PBM',
    # Camera RAW formats (limited support)
    'cr2': 'RAW',
    'cr3': 'RAW',
    'nef': 'RAW',
    'arw': 'RAW',
    'dng': 'RAW',
    # Other formats
    'svg': 'SVG',
    'heic': 'HEIC',
    'heif': 'HEIF',
}

# PDF Page size options
PDF_PAGE_SIZES = {
    'auto': None,  # Auto-detect from image
    'letter': (612, 792),  # 8.5 x 11 inches in points
    'a4': (595.28, 841.89),  # A4 in points
    'legal': (612, 1008),  # 8.5 x 14 inches
    'tabloid': (792, 1224),  # 11 x 17 inches
    'a3': (841.89, 1190.55),  # A3 in points
    'a5': (419.53, 595.28),  # A5 in points
}

# Quality presets (DPI for output)
QUALITY_PRESETS = {
    'low': 72,      # 72 DPI - web quality
    'medium': 150,  # 150 DPI - screen quality
    'high': 300,    # 300 DPI - print quality
    'maximum': 600, # 600 DPI - high quality print
}


def is_supported_image(filename: str) -> bool:
    """Check if file extension is a supported image format"""
    ext = Path(filename).suffix.lower().replace('.', '')
    return ext in SUPPORTED_IMAGE_FORMATS


def get_image_format_name(filename: str) -> str:
    """Get the format name for saving images"""
    ext = Path(filename).suffix.lower().replace('.', '')
    if ext == 'jpg':
        ext = 'jpeg'
    return SUPPORTED_IMAGE_FORMATS.get(ext, ext.upper())


def convert_image_to_pdf(
    image_path: Path,
    page_size: str = 'auto',
    quality: str = 'high',
    margin: float = 0
) -> Path:
    """Convert a single image to PDF with configurable options.
    
    Args:
        image_path: Path to the input image
        page_size: Page size ('auto', 'letter', 'a4', 'legal', 'tabloid', 'a3', 'a5')
        quality: Quality preset ('low', 'medium', 'high', 'maximum')
        margin: Margin in points (default 0)
    
    Returns:
        Path to the generated PDF
    """
    img = Image.open(image_path)
    
    # Get image dimensions in pixels
    img_width_px, img_height_px = img.size
    
    # Convert to RGB if necessary (required for PDF)
    if img.mode in ('RGBA', 'LA', 'P'):
        img = img.convert('RGB')
    
    # Determine output page size
    if page_size == 'auto':
        # Use image aspect ratio with a reasonable max size
        # Convert pixels to points (72 DPI baseline)
        max_width = 612  # Letter width in points
        scale = max_width / img_width_px
        page_width = img_width_px * scale
        page_height = img_height_px * scale
    else:
        page_size_tuple = PDF_PAGE_SIZES.get(page_size, PDF_PAGE_SIZES['letter'])
        page_width, page_height = page_size_tuple
    
    # Apply margin
    content_width = page_width - (2 * margin)
    content_height = page_height - (2 * margin)
    
    # Calculate scaling to fit image in page while maintaining aspect ratio
    scale_w = content_width / img_width_px if img_width_px > content_width else 1
    scale_h = content_height / img_height_px if img_height_px > content_height else 1
    scale = min(scale_w, scale_h)
    
    # Calculate final image dimensions on PDF
    final_width = img_width_px * scale
    final_height = img_height_px * scale
    
    # Calculate position to center the image
    left = (page_width - final_width) / 2
    top = (page_height - final_height) / 2
    
    # Create PDF with reportlab for better quality control
    from reportlab.lib.pagesizes import landscape, portrait
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import inch
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pdf"
    c = canvas.Canvas(str(output_path), pagesize=(page_width, page_height))
    
    # Set quality-based compression
    quality_dpi = QUALITY_PRESETS.get(quality, QUALITY_PRESETS['high'])
    
    # Draw image
    c.drawImage(
        str(image_path),
        left, top,
        width=final_width,
        height=final_height,
        preserveAspectRatio=True,
        mask='auto'
    )
    
    c.save()
    
    return output_path


def convert_multiple_images_to_pdf(
    image_paths: List[Path],
    page_size: str = 'auto',
    quality: str = 'high',
    margin: float = 0,
    one_image_per_page: bool = True
) -> Path:
    """Convert multiple images to a single PDF.
    
    Args:
        image_paths: List of paths to input images (in order)
        page_size: Page size for all pages
        quality: Quality preset
        margin: Margin in points
        one_image_per_page: If True, each image gets its own page
    
    Returns:
        Path to the generated PDF
    """
    from reportlab.lib.pagesizes import landscape, portrait
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import inch
    
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pdf"
    c = canvas.Canvas(str(output_path))
    
    for idx, image_path in enumerate(image_paths):
        if idx > 0:
            c.showPage()  # New page for each image (except first)
        
        img = Image.open(image_path)
        img_width_px, img_height_px = img.size
        
        # Convert to RGB if necessary
        if img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGB')
        
        # Determine page size
        if page_size == 'auto':
            # Use image dimensions as page size
            # Convert pixels to points (72 DPI baseline)
            page_width = img_width_px
            page_height = img_height_px
        else:
            page_size_tuple = PDF_PAGE_SIZES.get(page_size, PDF_PAGE_SIZES['letter'])
            page_width, page_height = page_size_tuple
        
        # Set page size
        c.setPageSize((page_width, page_height))
        
        # Apply margin
        content_width = page_width - (2 * margin)
        content_height = page_height - (2 * margin)
        
        # Calculate scaling to fit image
        scale_w = content_width / img_width_px if img_width_px > content_width else 1
        scale_h = content_height / img_height_px if img_height_px > content_height else 1
        scale = min(scale_w, scale_h)
        
        final_width = img_width_px * scale
        final_height = img_height_px * scale
        
        # Center image on page
        left = (page_width - final_width) / 2
        top = (page_height - final_height) / 2
        
        # Draw image
        c.drawImage(
            str(image_path),
            left, top,
            width=final_width,
            height=final_height,
            preserveAspectRatio=True,
            mask='auto'
        )
    
    c.save()
    
    return output_path


def convert_images_to_pdf_zip(
    image_paths: List[Path],
    page_size: str = 'auto',
    quality: str = 'high',
    margin: float = 0
) -> Path:
    """Convert images to individual PDFs and return as ZIP.
    
    Each image is converted to a separate PDF file.
    """
    pdf_paths = []
    
    for image_path in image_paths:
        try:
            pdf_path = convert_image_to_pdf(
                image_path,
                page_size=page_size,
                quality=quality,
                margin=margin
            )
            pdf_paths.append(pdf_path)
        except Exception as e:
            print(f"Failed to convert {image_path}: {e}")
            continue
    
    if not pdf_paths:
        raise Exception("No images could be converted")
    
    # Create ZIP with individual PDFs
    zip_name = f"images_to_pdf_{uuid.uuid4()}"
    zip_path = create_zip(pdf_paths, zip_name)
    
    # Clean up individual PDFs
    for pdf_path in pdf_paths:
        try:
            pdf_path.unlink()
        except:
            pass
    
    return zip_path

# ============== FastAPI Routes ==============


# ============== DOCX Conversions ==============

@api_router.post("/docx-to-pdf")
async def docx_to_pdf(
    file: UploadFile = File(...),
    target_format: str = Form("pdf")
):
    """Convert DOCX to PDF"""
    try:
        input_path = save_upload_file_tmp(file)
        output_path = convert_docx_to_pdf(input_path)

        # Save to history
        history = ConversionHistory(
            conversion_type="document",
            source_format="docx",
            target_format="pdf",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        await db.conversion_history.insert_one(doc)

        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".docx", ".pdf"),
            media_type="application/pdf"
        )
    except Exception as e:
        # Save failed history
        history = ConversionHistory(
            conversion_type="document",
            source_format="docx",
            target_format="pdf",
            filename=file.filename,
            status="failed"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=str(e))


@api_router.post("convert/docx-to-doc")
async def docx_to_doc(
    file: UploadFile = File(...),
    target_format: str = Form("doc")
):
    """Convert DOCX to DOC"""
    try:
        input_path = save_upload_file_tmp(file)
        output_path = convert_docx_to_doc(input_path)

        # Save to history
        history = ConversionHistory(
            conversion_type="document",
            source_format="docx",
            target_format="doc",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        await db.conversion_history.insert_one(doc)

        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".docx", ".doc"),
            media_type="application/msword"
        )
    except Exception as e:
        # Save failed history
        history = ConversionHistory(
            conversion_type="document",
            source_format="docx",
            target_format="doc",
            filename=file.filename,
            status="failed"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=str(e))


@api_router.post("/docx-to-txt")
async def docx_to_txt(
    file: UploadFile = File(...),
    target_format: str = Form("txt")
):
    """Convert DOCX to Text"""
    input_path = save_upload_file_tmp(file)
    output_path = convert_docx_to_txt(input_path)
    
    # await save_conversion_history("document", "docx", "txt", file.filename)
    
    return FileResponse(
        path=output_path,
        filename="converted.txt",
        media_type="text/plain"
    )


# ============== DOC Conversions ==============

@api_router.post("/doc-to-docx")
async def doc_to_docx(
    file: UploadFile = File(...),
    target_format: str = Form("docx")
):
    """Convert DOC to DOCX"""
    input_path = save_upload_file_tmp(file)
    output_path = convert_doc_to_docx(input_path)
    
    # await save_conversion_history("document", "doc", "docx", file.filename)
    
    return FileResponse(
        path=output_path,
        filename="converted.docx",
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )


@api_router.post("/doc-to-pdf")
async def doc_to_pdf(
    file: UploadFile = File(...),
    target_format: str = Form("pdf")
):
    """Convert DOC to PDF (via DOCX)"""
    input_path = save_upload_file_tmp(file)
    output_path = convert_doc_to_pdf(input_path)
    
    # await save_conversion_history("document", "doc", "pdf", file.filename)
    
    return FileResponse(
        path=output_path,
        filename="converted.pdf",
        media_type="application/pdf"
    )


# ============== PDF Conversions ==============

@api_router.post("/convert/pdf-to-docx")
async def pdf_to_docx(
    file: UploadFile = File(...),
    target_format: str = Form("docx")
):
    """Convert PDF to DOCX"""
    try:
        # 1️⃣ Save uploaded PDF
        input_path = save_upload_file_tmp(file)

        # 2️⃣ Convert PDF → DOCX
        output_path = convert_pdf_to_docx(input_path)

        # 3️⃣ SAVE HISTORY (SUCCESS)
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="docx",
            filename=file.filename,
            status="success"
        )

        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()

        await db.conversion_history.insert_one(doc)

        # 4️⃣ Return converted file
        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".pdf", ".docx"),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    except HTTPException:
        raise

    except Exception as e:
        # 5️⃣ SAVE HISTORY (FAILED)
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="docx",
            filename=file.filename,
            status="failed"
        )

        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)

        await db.conversion_history.insert_one(doc)

        raise HTTPException(
            status_code=500,
            detail=f"PDF to DOCX conversion failed: {str(e)}"
        )


@api_router.post("/convert/pdf-to-doc")
async def pdf_to_doc(
    file: UploadFile = File(...),
    target_format: str = Form("doc")
):
    """Convert PDF to DOC (via DOCX)"""
    try:
        input_path = save_upload_file_tmp(file)
        output_path = convert_pdf_to_doc(input_path)

        # Save to history
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="doc",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        await db.conversion_history.insert_one(doc)

        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".pdf", ".doc"),
            media_type="application/msword"
        )
    except Exception as e:
        # Save failed history
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="doc",
            filename=file.filename,
            status="failed"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=str(e))


@api_router.post("/pdf-to-txt")
async def pdf_to_txt(
    file: UploadFile = File(...),
    target_format: str = Form("txt")
):
    """Convert PDF to Text"""
    try:
        print(f"Received PDF to Text conversion request: {file.filename}")
        
        # Validate file type
        if not file.filename.lower().endswith('.pdf'):
            raise HTTPException(
                status_code=400,
                detail="Invalid file type. Please upload a PDF file."
            )
        
        # Save uploaded PDF
        input_path = save_upload_file_tmp(file)
        print(f"Saved temporary file: {input_path}")
        
        # Convert PDF to text
        output_path = convert_pdf_to_text(input_path)
        print(f"Converted to text: {output_path}")
        
        # Verify output file exists and has content
        if not output_path.exists() or output_path.stat().st_size == 0:
            raise HTTPException(
                status_code=500,
                detail="Text extraction failed. The PDF may not contain extractable text."
            )

        # Read extracted text for logging
        extracted_text = output_path.read_text()
        print(f"Extracted {len(extracted_text)} characters from PDF")

        # Save to history
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="txt",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        await db.conversion_history.insert_one(doc)

        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".pdf", ".txt"),
            media_type="text/plain"
        )
    except HTTPException:
        raise
    except Exception as e:
        print(f"PDF to Text conversion error: {str(e)}")
        # Save failed history
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="txt",
            filename=file.filename if 'file' in locals() else "unknown",
            status="failed"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)
        try:
            await db.conversion_history.insert_one(doc)
        except:
            pass
        raise HTTPException(
            status_code=500,
            detail=f"PDF to Text conversion failed: {str(e)}"
        )


@api_router.post("convert/pdf-to-xlsx")
async def pdf_to_xlsx(
    file: UploadFile = File(...),
    target_format: str = Form("xlsx")
):
    """Convert PDF to Excel"""
    try:
        input_path = save_upload_file_tmp(file)
        output_path = convert_pdf_to_excel(input_path)

        # Save to history
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="xlsx",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        await db.conversion_history.insert_one(doc)

        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".pdf", ".xlsx"),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
    except Exception as e:
        # Save failed history
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="xlsx",
            filename=file.filename,
            status="failed"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=str(e))


@api_router.post("/pdf-to-pptx")
async def pdf_to_pptx(
    file: UploadFile = File(...),
    target_format: str = Form("pptx")
):
    """Convert PDF to PowerPoint"""
    try:
        input_path = save_upload_file_tmp(file)
        output_path = convert_pdf_to_pptx(input_path)

        # Save to history
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="pptx",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        await db.conversion_history.insert_one(doc)

        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".pdf", ".pptx"),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        # Save failed history
        history = ConversionHistory(
            conversion_type="document",
            source_format="pdf",
            target_format="pptx",
            filename=file.filename,
            status="failed"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=str(e))
    


@api_router.post("convert/text-to-pdf")
async def text_to_pdf(
    file: UploadFile = File(...),
    target_format: str = Form("pdf")
):
    """Convert Text to PDF"""
    try:
        # 1️⃣ Save uploaded text file
        input_path = save_upload_file_tmp(file)

        # 2️⃣ Convert TXT → PDF
        output_path = convert_text_to_pdf(input_path)

        # 3️⃣ Save to history (SUCCESS)
        history = ConversionHistory(
            conversion_type="document",
            source_format="txt",
            target_format="pdf",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        await db.conversion_history.insert_one(doc)

        # 4️⃣ Return PDF
        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".txt", ".pdf"),
            media_type="application/pdf"
        )

    except Exception as e:
        # 5️⃣ Save to history (FAILED)
        history = ConversionHistory(
            conversion_type="document",
            source_format="txt",
            target_format="pdf",
            filename=file.filename,
            status="failed"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)
        await db.conversion_history.insert_one(doc)

        raise HTTPException(
            status_code=500,
            detail=str(e)
        )


# ============== Excel Conversions ==============

@api_router.post("/xlsx-to-pdf")
async def xlsx_to_pdf(
    file: UploadFile = File(...),
    target_format: str = Form("pdf")
):
    """Convert Excel XLSX to PDF"""
    try:
        input_path = save_upload_file_tmp(file)
        output_path = convert_excel_to_pdf(input_path)

        # Save to history (success)
        history = ConversionHistory(
            conversion_type="document",
            source_format="xlsx",
            target_format="pdf",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        await db.conversion_history.insert_one(doc)

        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".xlsx", ".pdf"),
            media_type="application/pdf"
        )
    except Exception as e:
        # Save failed history
        history = ConversionHistory(
            conversion_type="document",
            source_format="xlsx",
            target_format="pdf",
            filename=file.filename,
            status="failed"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=str(e))


@api_router.post("/xls-to-pdf")
async def xls_to_pdf(
    file: UploadFile = File(...),
    target_format: str = Form("pdf")
):
    """Convert Excel XLS to PDF"""
    try:
        input_path = save_upload_file_tmp(file)
        output_path = convert_excel_to_pdf(input_path)

        # Save to history (success)
        history = ConversionHistory(
            conversion_type="document",
            source_format="xls",
            target_format="pdf",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        await db.conversion_history.insert_one(doc)

        return FileResponse(
            path=output_path,
            filename=file.filename.replace(".xls", ".pdf"),
            media_type="application/pdf"
        )
    except Exception as e:
        # Save failed history
        history = ConversionHistory(
            conversion_type="document",
            source_format="xls",
            target_format="pdf",
            filename=file.filename,
            status="failed"
        )
        doc = history.model_dump()
        doc["timestamp"] = doc["timestamp"].isoformat()
        doc["error"] = str(e)
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=str(e))


# ============== PowerPoint Conversions ==============

@api_router.post("/pptx-to-pdf")
async def pptx_to_pdf(
    file: UploadFile = File(...),
    target_format: str = Form("pdf")
):
    """Convert PowerPoint PPTX to PDF"""
    input_path = save_upload_file_tmp(file)
    output_path = convert_pptx_to_pdf(input_path)
    
    # await save_conversion_history("document", "pptx", "pdf", file.filename)
    
    return FileResponse(
        path=output_path,
        filename="converted.pdf",
        media_type="application/pdf"
    )

@api_router.post("/ppt-to-pdf")
async def ppt_to_pdf(
    file: UploadFile = File(...),
    target_format: str = Form("pdf")
):
    input_path = save_upload_file_tmp(file)

    try:
        output_path = convert_ppt_to_pdf(input_path)
    finally:
        input_path.unlink(missing_ok=True)

    return FileResponse(
        path=str(output_path),
        media_type="application/pdf",
        filename=f"{Path(file.filename).stem}.pdf"
    )

# ============== Text Conversions ==============

@api_router.post("/txt-to-docx")
async def txt_to_docx(
    file: UploadFile = File(...),
    target_format: str = Form("docx")
):
    """Convert Text to DOCX"""
    input_path = save_upload_file_tmp(file)
    output_path = convert_text_to_docx(input_path)
    
    # await save_conversion_history("document", "txt", "docx", file.filename)
    
    return FileResponse(
        path=output_path,
        filename="converted.docx",
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )


@api_router.post("/txt-to-pdf")
async def txt_to_pdf(
    file: UploadFile = File(...),
    target_format: str = Form("pdf")
):
    """Convert Text to PDF"""
    input_path = save_upload_file_tmp(file)
    output_path = convert_text_to_pdf(input_path)
    
    # await save_conversion_history("document", "txt", "pdf", file.filename)
    
    return FileResponse(
        path=output_path,
        filename="converted.pdf",
        media_type="application/pdf"
    )


# ============== Image Conversions ==============

@api_router.post("/image-to-pdf")
async def image_to_pdf(
    file: UploadFile = File(...),
    target_format: str = Form("pdf")
):
    """Convert Image (JPG, PNG) to PDF"""
    input_path = save_upload_file_tmp(file)
    output_path = convert_image_to_pdf(input_path)

    # Save to history
    history = ConversionHistory(
        conversion_type="image",
        source_format=input_path.suffix.lower().replace('.', ''),
        target_format="pdf",
        filename=file.filename,
        status="success"
    )
    doc = history.model_dump()
    doc['timestamp'] = doc['timestamp'].isoformat()
    await db.conversion_history.insert_one(doc)

    return FileResponse(
        path=output_path,
        filename="converted.pdf",
        media_type="application/pdf"
    )


@api_router.post("/images-to-pdf")
async def images_to_pdf(
    files: List[UploadFile] = File(...),
    page_size: str = Form("auto"),
    quality: str = Form("high"),
    margin: float = Form(0)
):
    """Convert multiple images to a single PDF.
    
    Each image becomes one page in the PDF.
    Supports JPG, PNG, WEBP, BMP, TIFF, GIF, and more.
    """
    try:
        if not files:
            raise HTTPException(status_code=400, detail="No files uploaded")
        
        if len(files) > 50:
            raise HTTPException(status_code=400, detail="Maximum 50 images allowed")
        
        # Validate page_size
        if page_size not in PDF_PAGE_SIZES:
            page_size = 'auto'
        
        # Validate quality
        if quality not in QUALITY_PRESETS:
            quality = 'high'
        
        # Save all uploaded files
        image_paths = []
        for file in files:
            if not is_supported_image(file.filename):
                continue
            try:
                input_path = save_upload_file_tmp(file)
                image_paths.append(input_path)
            except Exception as e:
                print(f"Failed to save {file.filename}: {e}")
                continue
        
        if not image_paths:
            raise HTTPException(status_code=400, detail="No supported image files found")
        
        # Convert images to PDF
        output_path = convert_multiple_images_to_pdf(
            image_paths,
            page_size=page_size,
            quality=quality,
            margin=margin
        )
        
        # Generate filename from first image
        first_image_name = files[0].filename
        base_name = Path(first_image_name).stem
        
        # Save to history
        history = ConversionHistory(
            conversion_type="image",
            source_format="multiple",
            target_format="pdf",
            filename=f"{len(files)} images",
            status="success"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.conversion_history.insert_one(doc)
        
        return FileResponse(
            path=output_path,
            filename=f"{base_name}_and_{len(files)-1}_more.pdf",
            media_type="application/pdf"
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@api_router.post("/images-to-pdf-individual")
async def images_to_pdf_individual(
    files: List[UploadFile] = File(...),
    page_size: str = Form("auto"),
    quality: str = Form("high"),
    margin: float = Form(0)
):
    """Convert multiple images to individual PDFs and return as ZIP.
    
    Each image is converted to a separate PDF file.
    """
    try:
        if not files:
            raise HTTPException(status_code=400, detail="No files uploaded")
        
        if len(files) > 50:
            raise HTTPException(status_code=400, detail="Maximum 50 images allowed")
        
        # Validate page_size
        if page_size not in PDF_PAGE_SIZES:
            page_size = 'auto'
        
        # Validate quality
        if quality not in QUALITY_PRESETS:
            quality = 'high'
        
        # Save all uploaded files
        image_paths = []
        filenames = []
        for file in files:
            if not is_supported_image(file.filename):
                continue
            try:
                input_path = save_upload_file_tmp(file)
                image_paths.append(input_path)
                filenames.append(Path(file.filename).stem)
            except Exception as e:
                print(f"Failed to save {file.filename}: {e}")
                continue
        
        if not image_paths:
            raise HTTPException(status_code=400, detail="No supported image files found")
        
        # Convert images to individual PDFs
        output_path = convert_images_to_pdf_zip(
            image_paths,
            page_size=page_size,
            quality=quality,
            margin=margin
        )
        
        # Save to history
        history = ConversionHistory(
            conversion_type="image",
            source_format="multiple",
            target_format="pdf-zip",
            filename=f"{len(files)} images to individual PDFs",
            status="success"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.conversion_history.insert_one(doc)
        
        return FileResponse(
            path=output_path,
            filename="individual_pdfs.zip",
            media_type="application/zip"
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@api_router.get("/images/formats")
async def get_supported_image_formats():
    """Get list of supported image formats for Image to PDF conversion"""
    formats = []
    for ext, name in SUPPORTED_IMAGE_FORMATS.items():
        formats.append({
            "extension": f".{ext}",
            "name": name,
            "supported": True
        })
    return {
        "formats": formats,
        "page_sizes": list(PDF_PAGE_SIZES.keys()),
        "quality_presets": list(QUALITY_PRESETS.keys())
    }


# ============== Generic Document Conversion ==============

@api_router.post("/convert/document")
async def convert_document(
    file: UploadFile = File(...),
    target_format: str = Form(...)
):
    """Generic document conversion endpoint"""
    try:
        input_path = save_upload_file_tmp(file)
        source_format = input_path.suffix.lower().replace('.', '')

        # Define conversion function mapping
        conversion_functions = {
            ("pdf", "docx"): convert_pdf_to_docx,
            ("pdf", "doc"): convert_pdf_to_doc,
            ("pdf", "txt"): convert_pdf_to_text,
            ("pdf", "xlsx"): convert_pdf_to_excel,
            ("pdf", "pptx"): convert_pdf_to_pptx,
            ("docx", "pdf"): convert_docx_to_pdf,
            # ("docx", "doc"): convert_docx_to_doc,  # Disabled - function not implemented
            # ("docx", "txt"): convert_docx_to_txt,  # Disabled - function not implemented
            # ("doc", "docx"): convert_doc_to_docx,  # Disabled - function not implemented
            ("doc", "pdf"): convert_doc_to_pdf,
            ("xlsx", "pdf"): convert_excel_to_pdf,
            ("xls", "pdf"): convert_excel_to_pdf,
            ("pptx", "pdf"): convert_pptx_to_pdf,
            ("ppt", "pdf"): convert_ppt_to_pdf,
            ("txt", "docx"): convert_text_to_docx,
            ("txt", "pdf"): convert_text_to_pdf,
            ("jpg", "pdf"): convert_image_to_pdf,
            ("jpeg", "pdf"): convert_image_to_pdf,
            ("png", "pdf"): convert_image_to_pdf,
            ("webp", "pdf"): convert_image_to_pdf,
            ("bmp", "pdf"): convert_image_to_pdf,
        }

        # Normalize formats
        normalized_source = source_format.lower()
        if normalized_source == "jpg":
            normalized_source = "jpeg"

        normalized_target = target_format.lower()
        if normalized_target == "jpg":
            normalized_target = "jpeg"

        # Find conversion function
        conversion_key = (normalized_source, normalized_target)
        if conversion_key not in conversion_functions:
            raise HTTPException(
                status_code=400,
                detail=f"Conversion from {source_format} to {target_format} is not supported"
            )

        conversion_func = conversion_functions[conversion_key]
        output_path = conversion_func(input_path)

        # Define media types
        media_types = {
            "pdf": "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "doc": "application/msword",
            "txt": "text/plain",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "jpeg": "image/jpeg",
            "png": "image/png",
            "webp": "image/webp",
            "bmp": "image/bmp",
        }

        media_type = media_types.get(normalized_target, "application/octet-stream")

        # Save to history
        # await save_conversion_history("document", source_format, target_format, file.filename)

        return FileResponse(
            path=output_path,
            filename=f"converted.{target_format.lower()}",
            media_type=media_type
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# API Routes
@api_router.get("/")
async def root():
    return {"message": "File Conversion API"}

@api_router.get("/health")
async def health_check():
    """Health check endpoint for Docker container monitoring"""
    return {"status": "healthy", "service": "file-conversion-api"}

@api_router.post("/convert/image")
async def convert_image(
    file: UploadFile = File(...),
    target_format: str = Form(...)
):
    """Convert image between formats (JPG, PNG, WEBP, BMP)"""
    try:
        input_path = save_upload_file_tmp(file)
        source_format = input_path.suffix.lower().replace('.', '')

        # 🔧 FIX: Normalize JPG → JPEG for Pillow
        normalized_format = target_format.lower()
        if normalized_format == "jpg":
            normalized_format = "jpeg"

        output_path = convert_image_format(input_path, normalized_format)

        # Save to history
        history = ConversionHistory(
            conversion_type="image",
            source_format=source_format,
            target_format=target_format.lower(),
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.conversion_history.insert_one(doc)

        return FileResponse(
            path=output_path,
            filename=f"converted.{target_format.lower()}",
            media_type="application/octet-stream"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@api_router.post("/image/resize")
async def resize_image_endpoint(
    file: UploadFile = File(...),
    target_width: int = Form(...),
    target_height: int = Form(...),
    maintain_aspect_ratio: bool = Form(True),
    output_format: str = Form("jpeg"),
    quality: str = Form("high")
):
    """Resize an image to specified dimensions.
    
    Args:
        file: The image file to resize
        target_width: Target width in pixels
        target_height: Target height in pixels
        maintain_aspect_ratio: If True, maintains aspect ratio (default: True)
        output_format: Output format (jpeg, png, webp, bmp)
        quality: Quality preset (low, medium, high, maximum)
    """
    try:
        # Validate inputs
        if target_width < 1 or target_width > 10000:
            raise HTTPException(status_code=400, detail="Width must be between 1 and 10000 pixels")
        
        if target_height < 1 or target_height > 10000:
            raise HTTPException(status_code=400, detail="Height must be between 1 and 10000 pixels")
        
        # Validate output format
        valid_formats = ['jpeg', 'jpg', 'png', 'webp', 'bmp']
        if output_format.lower() not in valid_formats:
            raise HTTPException(status_code=400, detail=f"Invalid output format. Must be one of: {', '.join(valid_formats)}")
        
        # Normalize JPG to JPEG
        normalized_format = output_format.lower()
        if normalized_format == "jpg":
            normalized_format = "jpeg"
        
        # Validate quality
        valid_qualities = ['low', 'medium', 'high', 'maximum']
        if quality.lower() not in valid_qualities:
            quality = 'high'
        
        # Save uploaded file
        input_path = save_upload_file_tmp(file)
        
        # Get original dimensions for response
        with Image.open(input_path) as img:
            original_width, original_height = img.size
        
        # Resize the image
        output_path = resize_image(
            image_path=input_path,
            target_width=target_width,
            target_height=target_height,
            maintain_aspect_ratio=maintain_aspect_ratio,
            output_format=normalized_format,
            quality=quality.lower()
        )
        
        # Get new dimensions
        with Image.open(output_path) as resized_img:
            new_width, new_height = resized_img.size
        
        # Save to history
        history = ConversionHistory(
            conversion_type="image_resize",
            source_format=input_path.suffix.lower().replace('.', ''),
            target_format=normalized_format,
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.conversion_history.insert_one(doc)
        
        # Determine media type
        media_types = {
            'jpeg': 'image/jpeg',
            'jpg': 'image/jpeg',
            'png': 'image/png',
            'webp': 'image/webp',
            'bmp': 'image/bmp',
        }
        media_type = media_types.get(normalized_format, 'application/octet-stream')
        
        # Generate output filename
        original_name = Path(file.filename).stem
        output_filename = f"{original_name}_resized.{normalized_format}"
        
        return FileResponse(
            path=output_path,
            filename=output_filename,
            media_type=media_type
        )
    except HTTPException:
        raise
    except Exception as e:
        # Save failed history
        history = ConversionHistory(
            conversion_type="image_resize",
            source_format="unknown",
            target_format=output_format.lower(),
            filename=file.filename if 'file' in locals() else "unknown",
            status="failed"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        doc['error'] = str(e)
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=str(e))


@api_router.post("/images/resize")
async def resize_images_endpoint(
    files: List[UploadFile] = File(...),
    target_width: int = Form(...),
    target_height: int = Form(...),
    maintain_aspect_ratio: bool = Form(True),
    output_format: str = Form("jpeg"),
    quality: str = Form("high")
):
    """Resize multiple images and return as ZIP.
    
    Args:
        files: List of image files to resize
        target_width: Target width in pixels
        target_height: Target height in pixels
        maintain_aspect_ratio: If True, maintains aspect ratio
        output_format: Output format (jpeg, png, webp, bmp)
        quality: Quality preset (low, medium, high, maximum)
    """
    try:
        if not files:
            raise HTTPException(status_code=400, detail="No files uploaded")
        
        if len(files) > 20:
            raise HTTPException(status_code=400, detail="Maximum 20 images allowed for batch resize")
        
        # Validate inputs
        if target_width < 1 or target_width > 10000:
            raise HTTPException(status_code=400, detail="Width must be between 1 and 10000 pixels")
        
        if target_height < 1 or target_height > 10000:
            raise HTTPException(status_code=400, detail="Height must be between 1 and 10000 pixels")
        
        # Validate output format
        valid_formats = ['jpeg', 'jpg', 'png', 'webp', 'bmp']
        if output_format.lower() not in valid_formats:
            raise HTTPException(status_code=400, detail=f"Invalid output format. Must be one of: {', '.join(valid_formats)}")
        
        # Normalize JPG to JPEG
        normalized_format = output_format.lower()
        if normalized_format == "jpg":
            normalized_format = "jpeg"
        
        # Validate quality
        valid_qualities = ['low', 'medium', 'high', 'maximum']
        if quality.lower() not in valid_qualities:
            quality = 'high'
        
        # Save all uploaded files
        image_paths = []
        for file in files:
            if not is_supported_image(file.filename):
                continue
            try:
                input_path = save_upload_file_tmp(file)
                image_paths.append(input_path)
            except Exception as e:
                print(f"Failed to save {file.filename}: {e}")
                continue
        
        if not image_paths:
            raise HTTPException(status_code=400, detail="No supported image files found")
        
        # Resize images
        output_path = resize_multiple_images(
            image_paths=image_paths,
            target_width=target_width,
            target_height=target_height,
            maintain_aspect_ratio=maintain_aspect_ratio,
            output_format=normalized_format,
            quality=quality.lower()
        )
        
        # Save to history
        history = ConversionHistory(
            conversion_type="image_resize",
            source_format="multiple",
            target_format=f"{normalized_format}-zip",
            filename=f"{len(files)} images",
            status="success"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.conversion_history.insert_one(doc)
        
        return FileResponse(
            path=output_path,
            filename="resized_images.zip",
            media_type="application/zip"
        )
    except HTTPException:
        raise
@api_router.post("/pdf/lock")
async def lock_pdf_endpoint(
    file: UploadFile = File(...),
    password: str = Form(...)
):
    """Lock/encrypt PDF with password"""
    try:
        input_path = save_upload_file_tmp(file)
        output_path = lock_pdf(input_path, password)
        
        return FileResponse(
            path=output_path,
            filename="locked.pdf",
            media_type="application/pdf"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/pdf/unlock")
async def unlock_pdf_endpoint(
    file: UploadFile = File(...),
    password: str = Form(...)
):
    """Unlock/decrypt PDF with password"""
    try:
        input_path = save_upload_file_tmp(file)
        output_path = unlock_pdf(input_path, password)
        
        return FileResponse(
            path=output_path,
            filename="unlocked.pdf",
            media_type="application/pdf"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/pdf/merge")
async def merge_pdfs_endpoint(files: List[UploadFile] = File(...)):
    """Merge multiple PDFs"""
    try:
        # Validate file count
        MAX_FILES = 20
        MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
        
        if len(files) > MAX_FILES:
            raise HTTPException(
                status_code=400, 
                detail=f"Maximum {MAX_FILES} PDF files allowed for merge. You provided {len(files)} files."
            )
        
        if len(files) < 2:
            raise HTTPException(
                status_code=400,
                detail="At least 2 PDF files are required for merge"
            )
        
        valid_pdf_paths = []
        
        # Process each file
        for file in files:
            if file is None:
                continue
            
            filename = getattr(file, 'filename', None) or 'unknown'
            
            # Check if filename indicates PDF
            if not filename.lower().endswith('.pdf'):
                raise HTTPException(
                    status_code=400,
                    detail=f"Invalid file type: {filename}. Only PDF files are allowed for merge."
                )
            
            # Save to temp directory (this handles reading the file content)
            try:
                pdf_path = save_upload_file_tmp(file)
                
                # Validate the saved file
                if not pdf_path.exists():
                    raise HTTPException(
                        status_code=500,
                        detail=f"Failed to save file: {filename}"
                    )
                
                file_size = pdf_path.stat().st_size
                
                if file_size == 0:
                    raise HTTPException(
                        status_code=400,
                        detail=f"File is empty: {filename}"
                    )
                
                if file_size > MAX_FILE_SIZE:
                    raise HTTPException(
                        status_code=400,
                        detail=f"File too large: {filename}. Maximum size is 50MB."
                    )
                
                # Verify it's a valid PDF by checking file header
                try:
                    with open(pdf_path, 'rb') as f:
                        header = f.read(5)
                        if header[:5] != b'%PDF-':
                            raise HTTPException(
                                status_code=400,
                                detail=f"Invalid PDF file: {filename}"
                            )
                except Exception as pdf_error:
                    raise HTTPException(
                        status_code=400,
                        detail=f"Invalid PDF file: {filename}"
                    )
                
                valid_pdf_paths.append(pdf_path)
                
            except HTTPException:
                raise
            except Exception as save_error:
                raise HTTPException(
                    status_code=500,
                    detail=f"Error processing file {filename}: {str(save_error)}"
                )
        
        if len(valid_pdf_paths) < 2:
            raise HTTPException(
                status_code=400,
                detail="At least 2 valid PDF files are required for merge"
            )
        
        output_path = merge_pdfs(valid_pdf_paths)
        
        return FileResponse(
            path=output_path,
            filename="merged.pdf",
            media_type="application/pdf"
        )
    except HTTPException:
        raise
    except Exception as e:
        # Log the error for debugging
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/pdf/split")
async def split_pdf_endpoint(
    file: UploadFile = File(...),
    page_ranges: str = Form(...)
):
    """Split PDF (e.g., page_ranges='1-3,4-6' or '1,3,5')"""
    try:
        input_path = save_upload_file_tmp(file)
        output_paths = split_pdf(input_path, page_ranges)
        
        # Create ZIP with all split PDFs
        zip_path = create_zip(output_paths, "split_pdfs")
        
        return FileResponse(
            path=zip_path,
            filename="split_pdfs.zip",
            media_type="application/zip"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/zip/compress")
async def compress_files(files: List[UploadFile] = File(...)):
    """Compress multiple files into ZIP"""
    try:
        file_paths = [save_upload_file_tmp(file) for file in files]
        zip_path = create_zip(file_paths, "compressed_files")
        
        return FileResponse(
            path=zip_path,
            filename="compressed.zip",
            media_type="application/zip"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/zip/extract")
async def extract_files(file: UploadFile = File(...)):
    """Extract ZIP archive and return list of files for individual download"""
    try:
        input_path = save_upload_file_tmp(file)

        # Validate that the file is a valid ZIP
        if not zipfile.is_zipfile(input_path):
            raise HTTPException(status_code=400, detail="Invalid ZIP file")

        # Create extraction directory
        extracted_dir = TEMP_DIR / f"{uuid.uuid4()}_extracted"
        extracted_dir.mkdir(exist_ok=True)

        # Extract ZIP preserving structure
        try:
            with zipfile.ZipFile(input_path, 'r') as zipf:
                zipf.extractall(extracted_dir)
        except zipfile.BadZipFile:
            raise HTTPException(status_code=400, detail="Corrupt ZIP file - unable to extract")

        # Get all extracted files preserving folder structure
        extracted_paths = list(extracted_dir.glob('**/*'))
        extracted_files = [p for p in extracted_paths if p.is_file()]

        if len(extracted_files) == 0:
            raise HTTPException(status_code=400, detail="ZIP file is empty")

        # Create file list with metadata
        file_list = []
        for file_path in extracted_files:
            relative_path = file_path.relative_to(extracted_dir)
            file_list.append({
                "path": str(relative_path),
                "size": file_path.stat().st_size,
                "extraction_id": extracted_dir.name
            })

        return {
            "files": file_list,
            "extraction_id": extracted_dir.name,
            "total_files": len(file_list)
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/zip/download-file/{extraction_id}/{file_path:path}")
async def download_file(extraction_id: str, file_path: str):
    """Download individual file from extracted ZIP"""
    try:
        extracted_dir = TEMP_DIR / extraction_id

        if not extracted_dir.exists():
            raise HTTPException(status_code=404, detail="Extraction session expired")

        file_path_obj = extracted_dir / file_path

        if not file_path_obj.exists() or not file_path_obj.is_file():
            raise HTTPException(status_code=404, detail="File not found")

        return FileResponse(
            path=file_path_obj,
            filename=file_path_obj.name,
            media_type="application/octet-stream"
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/ocr/languages")
async def get_ocr_languages():
    """Get list of available OCR languages with their names"""
    # Return all languages from LANGUAGE_NAMES dictionary
    # This includes both installed and available Tesseract languages
    languages = []
    
    # First, get available Tesseract languages
    try:
        available_langs = get_tesseract_languages()
        global AVAILABLE_OCR_LANGUAGES
        AVAILABLE_OCR_LANGUAGES = available_langs
    except Exception as e:
        print(f"Failed to refresh Tesseract languages: {e}")
        available_langs = AVAILABLE_OCR_LANGUAGES
    
    # Build response with all languages from LANGUAGE_NAMES
    for lang_code, lang_name in LANGUAGE_NAMES.items():
        if isinstance(lang_name, str):
            languages.append({
                "code": lang_code,
                "name": lang_name,
                "installed": lang_code in available_langs or lang_code.lower() in [l.lower() for l in available_langs]
            })
    
    # Sort by name for better UX
    languages.sort(key=lambda x: x['name'])
    
    return {
        "languages": languages,
        "count": len(languages)
    }

@api_router.post("/ocr/detect-language")
async def detect_ocr_language(file: UploadFile = File(...)):
    """Detect language/script from image using Tesseract OSD"""
    try:
        input_path = save_upload_file_tmp(file)
        result = detect_language_from_image(input_path)
        
        # Get full language info for suggested languages
        suggested_languages_info = []
        for lang_code in result["suggested_languages"]:
            lang_info = {
                "code": lang_code,
                "name": LANGUAGE_NAMES.get(lang_code, lang_code.title())
            }
            suggested_languages_info.append(lang_info)
        
        return {
            "detected_script": result["detected_script"],
            "orientation": result["orientation"],
            "confidence": result["confidence"],
            "suggested_languages": suggested_languages_info,
            "primary_language": suggested_languages_info[0] if suggested_languages_info else None
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Language detection failed: {str(e)}")

@api_router.post("/ocr/extract")
async def extract_text_ocr(
    file: UploadFile = File(...),
    language: str = Form("eng")
):
    """Extract text from image using OCR"""
    try:
        # Check if language is available
        if language not in AVAILABLE_OCR_LANGUAGES:
            available = ", ".join(AVAILABLE_OCR_LANGUAGES[:10])
            if len(AVAILABLE_OCR_LANGUAGES) > 10:
                available += f" and {len(AVAILABLE_OCR_LANGUAGES) - 10} more"
            raise HTTPException(
                status_code=400,
                detail=f"Language '{language}' is not installed. Available languages: {available}. Please install missing language packs for Tesseract."
            )

        input_path = save_upload_file_tmp(file)
        text = ocr_image(input_path, language)

        # Check if OCR returned an error message
        if text.startswith("OCR Error:"):
            raise HTTPException(status_code=500, detail=text)

        # Save to history
        history = ConversionHistory(
            conversion_type="ocr",
            source_format=input_path.suffix.lower().replace('.', ''),
            target_format="text",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.conversion_history.insert_one(doc)

        return {"text": text, "filename": file.filename, "language": language}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/search/pdf")
async def search_in_pdf_endpoint(
    file: UploadFile = File(...),
    search_term: str = Form(...)
):
    """Search for text within PDF document"""
    try:
        if not search_term or not search_term.strip():
            raise HTTPException(status_code=400, detail="Search term cannot be empty")

        input_path = save_upload_file_tmp(file)
        results = search_in_pdf(input_path, search_term.strip())

        # Save to history
        history = ConversionHistory(
            conversion_type="search",
            source_format="pdf",
            target_format="results",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.conversion_history.insert_one(doc)

        return results
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Search failed: {str(e)}")

@api_router.get("/history", response_model=List[ConversionHistory])
async def get_conversion_history():
    """Get conversion history"""
    history = await db.conversion_history.find({}, {"_id": 0}).sort("timestamp", -1).to_list(100)
    
    for item in history:
        if isinstance(item['timestamp'], str):
            item['timestamp'] = datetime.fromisoformat(item['timestamp'])
    
    return history

# ============== Watermark PDF Routes ==============

@api_router.post("/watermark/pdf/text")
async def add_text_watermark_endpoint(
    file: UploadFile = File(...),
    text: str = Form(...),
    font_name: str = Form("Helvetica-Bold"),
    font_size: int = Form(48),
    color: str = Form("#808080"),
    opacity: float = Form(0.3),
    rotation: float = Form(45),
    position: str = Form("center"),
    first_page_only: bool = Form(False),
    page_ranges: Optional[str] = Form(None),
    margin_x: float = Form(50),
    margin_y: float = Form(50),
    outline: bool = Form(False),
    outline_color: str = Form("#FFFFFF")
):
    """Add text watermark to PDF"""
    try:
        # Validate inputs
        if not text or not text.strip():
            raise HTTPException(status_code=400, detail="Watermark text cannot be empty")
        
        # Validate font_size
        if font_size < 8 or font_size > 200:
            raise HTTPException(status_code=400, detail="Font size must be between 8 and 200")
        
        # Validate opacity
        if opacity < 0 or opacity > 1:
            raise HTTPException(status_code=400, detail="Opacity must be between 0 and 1")
        
        # Validate rotation
        if rotation < -360 or rotation > 360:
            raise HTTPException(status_code=400, detail="Rotation must be between -360 and 360 degrees")
        
        # Validate position
        valid_positions = [POSITION_CENTER, POSITION_TOP_LEFT, POSITION_TOP_RIGHT, 
                          POSITION_BOTTOM_LEFT, POSITION_BOTTOM_RIGHT, POSITION_TILED]
        if position not in valid_positions:
            raise HTTPException(status_code=400, detail=f"Invalid position. Must be one of: {', '.join(valid_positions)}")
        
        # Log watermark request
        print(f"[WATERMARK] Processing text watermark request for file: {file.filename}")
        print(f"[WATERMARK] Text: {text}, Font: {font_name}, Size: {font_size}, Color: {color}")
        
        # Save uploaded PDF
        input_path = save_upload_file_tmp(file)
        print(f"[WATERMARK] Saved input file to: {input_path}")
        
        # Add watermark
        output_path = add_text_watermark(
            pdf_path=input_path,
            text=text,
            font_name=font_name,
            font_size=font_size,
            color=color,
            opacity=opacity,
            rotation=rotation,
            position=position,
            first_page_only=first_page_only,
            page_ranges=page_ranges,
            margin_x=margin_x,
            margin_y=margin_y,
            outline=outline,
            outline_color=outline_color
        )
        print(f"[WATERMARK] Generated watermarked PDF: {output_path}")
        
        # Save to history
        history = ConversionHistory(
            conversion_type="watermark",
            source_format="pdf",
            target_format="pdf",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.conversion_history.insert_one(doc)
        
        return FileResponse(
            path=output_path,
            filename=f"{Path(file.filename).stem}_watermarked.pdf",
            media_type="application/pdf"
        )
    except HTTPException:
        raise
    except Exception as e:
        # Enhanced error logging
        error_msg = f"Text watermark failed: {str(e)}"
        print(f"[WATERMARK ERROR] {error_msg}")
        import traceback
        traceback.print_exc()
        
        # Save failed history
        history = ConversionHistory(
            conversion_type="watermark",
            source_format="pdf",
            target_format="pdf",
            filename=file.filename if 'file' in locals() else "unknown",
            status="failed"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        doc['error'] = error_msg
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=error_msg)


@api_router.post("/watermark/pdf/image")
async def add_image_watermark_endpoint(
    file: UploadFile = File(...),
    watermark_file: UploadFile = File(...),
    opacity: float = Form(0.3),
    position: str = Form("center"),
    scale: float = Form(0.5),
    rotation: float = Form(0),
    first_page_only: bool = Form(False),
    page_ranges: Optional[str] = Form(None),
    margin_x: float = Form(50),
    margin_y: float = Form(50)
):
    """Add image/logo watermark to PDF"""
    try:
        # Validate opacity
        if opacity < 0 or opacity > 1:
            raise HTTPException(status_code=400, detail="Opacity must be between 0 and 1")
        
        # Validate scale
        if scale < 0.1 or scale > 2.0:
            raise HTTPException(status_code=400, detail="Scale must be between 0.1 and 2.0")
        
        # Validate rotation
        if rotation < -360 or rotation > 360:
            raise HTTPException(status_code=400, detail="Rotation must be between -360 and 360 degrees")
        
        # Validate position
        valid_positions = [POSITION_CENTER, POSITION_TOP_LEFT, POSITION_TOP_RIGHT, 
                          POSITION_BOTTOM_LEFT, POSITION_BOTTOM_RIGHT, POSITION_TILED]
        if position not in valid_positions:
            raise HTTPException(status_code=400, detail=f"Invalid position. Must be one of: {', '.join(valid_positions)}")
        
        # Log watermark request
        print(f"[WATERMARK] Processing image watermark request for file: {file.filename}")
        print(f"[WATERMARK] Watermark: {watermark_file.filename}, Opacity: {opacity}, Scale: {scale}")
        
        # Save uploaded PDF and watermark image
        input_path = save_upload_file_tmp(file)
        watermark_path = save_upload_file_tmp(watermark_file)
        print(f"[WATERMARK] Saved files - PDF: {input_path}, Image: {watermark_path}")
        
        # Add watermark
        output_path = add_image_watermark(
            pdf_path=input_path,
            image_path=watermark_path,
            opacity=opacity,
            position=position,
            scale=scale,
            rotation=rotation,
            first_page_only=first_page_only,
            page_ranges=page_ranges,
            margin_x=margin_x,
            margin_y=margin_y
        )
        print(f"[WATERMARK] Generated watermarked PDF: {output_path}")
        
        # Save to history
        history = ConversionHistory(
            conversion_type="watermark",
            source_format="pdf",
            target_format="pdf",
            filename=file.filename,
            status="success"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        await db.conversion_history.insert_one(doc)
        
        return FileResponse(
            path=output_path,
            filename=f"{Path(file.filename).stem}_watermarked.pdf",
            media_type="application/pdf"
        )
    except HTTPException:
        raise
    except Exception as e:
        # Enhanced error logging
        error_msg = f"Image watermark failed: {str(e)}"
        print(f"[WATERMARK ERROR] {error_msg}")
        import traceback
        traceback.print_exc()
        
        # Save failed history
        history = ConversionHistory(
            conversion_type="watermark",
            source_format="pdf",
            target_format="pdf",
            filename=file.filename if 'file' in locals() else "unknown",
            status="failed"
        )
        doc = history.model_dump()
        doc['timestamp'] = doc['timestamp'].isoformat()
        doc['error'] = error_msg
        await db.conversion_history.insert_one(doc)
        raise HTTPException(status_code=500, detail=error_msg)


# ============== Watermark Preview Endpoint ==============

@api_router.post("/watermark/pdf/preview")
async def watermark_preview_endpoint(
    file: UploadFile = File(...),
    watermark_type: str = Form("text"),  # "text" or "image"
    # Text watermark params
    text: Optional[str] = Form(None),
    font_name: str = Form("Helvetica-Bold"),
    font_size: int = Form(48),
    color: str = Form("#808080"),
    opacity: float = Form(0.3),
    rotation: float = Form(45),
    position: str = Form("center"),
    margin_x: float = Form(50),
    margin_y: float = Form(50),
    outline: bool = Form(False),
    outline_color: str = Form("#FFFFFF"),
    # Image watermark params
    watermark_file: Optional[UploadFile] = Form(None),
    scale: float = Form(0.5),
    # Common params
    page_number: int = Form(1)
):
    """Generate a preview of the watermark on a single PDF page.
    
    Returns a single-page PDF blob for real-time preview in the frontend.
    """
    try:
        # Validate inputs
        if watermark_type == "text":
            if not text or not text.strip():
                raise HTTPException(status_code=400, detail="Watermark text cannot be empty for text watermark")
        elif watermark_type == "image":
            if not watermark_file:
                raise HTTPException(status_code=400, detail="Watermark image is required for image watermark")
        else:
            raise HTTPException(status_code=400, detail="Invalid watermark type. Must be 'text' or 'image'")
        
        # Validate opacity
        if opacity < 0 or opacity > 1:
            raise HTTPException(status_code=400, detail="Opacity must be between 0 and 1")
        
        # Validate rotation
        if rotation < -360 or rotation > 360:
            raise HTTPException(status_code=400, detail="Rotation must be between -360 and 360 degrees")
        
        # Validate position
        valid_positions = [POSITION_CENTER, POSITION_TOP_LEFT, POSITION_TOP_RIGHT, 
                          POSITION_BOTTOM_LEFT, POSITION_BOTTOM_RIGHT, POSITION_TILED]
        if position not in valid_positions:
            raise HTTPException(status_code=400, detail=f"Invalid position. Must be one of: {', '.join(valid_positions)}")
        
        # Save uploaded PDF
        input_path = save_upload_file_tmp(file)
        
        if watermark_type == "text":
            # Add text watermark
            output_path = add_text_watermark(
                pdf_path=input_path,
                text=text,
                font_name=font_name,
                font_size=font_size,
                color=color,
                opacity=opacity,
                rotation=rotation,
                position=position,
                first_page_only=True,  # Only first page for preview
                page_ranges=None,
                margin_x=margin_x,
                margin_y=margin_y,
                outline=outline,
                outline_color=outline_color
            )
        else:
            # Save watermark image
            watermark_path = save_upload_file_tmp(watermark_file)
            
            # Add image watermark
            output_path = add_image_watermark(
                pdf_path=input_path,
                image_path=watermark_path,
                opacity=opacity,
                position=position,
                scale=scale,
                rotation=rotation,
                first_page_only=True,  # Only first page for preview
                page_ranges=None,
                margin_x=margin_x,
                margin_y=margin_y
            )
        
        print(f"[WATERMARK PREVIEW] Generated preview PDF: {output_path}")
        
        return FileResponse(
            path=output_path,
            filename="watermark_preview.pdf",
            media_type="application/pdf"
        )
    except HTTPException:
        raise
    except Exception as e:
        error_msg = f"Preview generation failed: {str(e)}"
        print(f"[WATERMARK PREVIEW ERROR] {error_msg}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=error_msg)


# ============== GET Available Options ==============

@api_router.get("/watermark/options")
async def get_watermark_options():
    """Get available watermark options"""
    return {
        "fonts": [
            "Helvetica",
            "Helvetica-Bold",
            "Helvetica-Oblique",
            "Helvetica-BoldOblique",
            "Times-Roman",
            "Times-Bold",
            "Times-Italic",
            "Times-BoldItalic",
            "Courier",
            "Courier-Bold",
            "Courier-Oblique",
            "Courier-BoldOblique",
            "Symbol",
            "ZapfDingbats"
        ],
        "positions": [
            {"value": "center", "label": "Center"},
            {"value": "top_left", "label": "Top Left"},
            {"value": "top_right", "label": "Top Right"},
            {"value": "bottom_left", "label": "Bottom Left"},
            {"value": "bottom_right", "label": "Bottom Right"},
            {"value": "tiled", "label": "Tiled (Repeated)"}
        ],
        "default_text_watermark": {
            "text": "CONFIDENTIAL",
            "font_name": "Helvetica-Bold",
            "font_size": 48,
            "color": "#808080",
            "opacity": 0.3,
            "rotation": 45,
            "position": "center"
        },
        "default_image_watermark": {
            "opacity": 0.3,
            "position": "center",
            "scale": 0.5,
            "rotation": 0
        }
    }


# Include the router in the main app
app.include_router(api_router)

# CORS configuration for HTTP-only access
def get_cors_origins():
    """Get CORS origins from environment or use sensible defaults for network access"""
    cors_env = os.environ.get('CORS_ORIGINS', '')
    if cors_env:
        origins = cors_env.split(',')
        # If wildcard is in the list, return just wildcard
        if '*' in origins:
            return ["*"]
        return origins
    
    # Check if running in Docker container
    in_docker = os.environ.get('REACT_APP_IN_DOCKER', '').lower() == 'true'
    
    if in_docker:
        # In Docker, allow all origins for internal network access
        # The Nginx reverse proxy handles security
        return ["*"]
    
    # Default origins for development (HTTP only)
    return [
        "http://localhost:3000",
        "http://localhost:3001",
        "http://localhost:3002",
        "http://127.0.0.1:3000",
        "http://127.0.0.1:3001",
        "http://127.0.0.1:3002",
        "*"  # Allow all origins for network access from different devices
    ]

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=get_cors_origins(),
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)

